#!/opt/venv/bin/python3
"""pptx_inspect - paginated structural inspection of .pptx decks.

Backed by python-pptx for slide/shape/placeholder/text traversal; the
shared `ooxml` helpers and stdlib `zipfile` are used for chart titles
and embedded media listing where python-pptx is awkward or silent.
"""

from __future__ import annotations

import argparse
import os
import re
import sys
import zipfile
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple

import ooxml
import pdf_text
import render
import render_publish
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from utils import (
    TEXT_PREVIEW_LIMIT,
    ellipsize,
    flatten_text,
    format_size,
    pad,
    parse_slide_patterns,
    safe_output,
)

from pptx_geometry import _text_extent_box, emu_to_inches, format_box, shape_kind

from pptx_typography import (
    _read_clr_map,
    _read_layout_chain,
    _read_theme_for_master,
    _resolve_scheme_color,
    format_placeholder_defaults,
    placeholder_type,
    resolve_placeholder_defaults,
    text_frame_lines,
)

from pptx_render_boxes import _annotate_boxes

import pptx_contrast
import pptx_grid
import pptx_validate

from pptx_audit import (
    BARE_MARGIN,
    BARE_RATE,
    LEFTOVER_LISTED,
    STRUCTURAL_KINDS,
    STRUCTURAL_LABELS,
    CoverRect,
    DENSITY_TOLERANCE,
    DeckFidelity,
    IMAGERY_RICH_TEMPLATE,
    IMAGERY_STRIP_RATIO,
    MIN_SLIDES_FOR_CONTENT_AUDIT,
    SHAPE_RETENTION_FLOOR,
    SlideContext,
    _count_slides,
    _cover_candidates,
    _deck_fidelity,
    _deck_structural_audit,
    _leading_break_audit,
    _drop_audit,
    _empty_table_rows_audit,
    _filler_audit,
    _hole_audit,
    _effective_font_size_pt,
    _fit_tokens,
    _has_embedded_blip,
    _is_leftover_suspect,
    _leftover_copy_audit,
    untouched_slides,
    _repeated_image_audit,
    _repeated_text_audit,
    _listed_slide_count,
    _package_names,
    _shape_text_iter,
    embedded_image,
    _shape_warning_markers,
    _split_sentence_markers,
    _void_markers,
    rendered_void,
    _slot_audit,
    slide_word_count,
)


DEFAULT_MAX_SHAPES = 200


USAGE = (
    "pptx_inspect <file> [--qa N [--boxes] | --slide N | --layouts | --text | "
    "--media | --render [--grid [--grid-cols N]] | --validate] [--compare FILE] "
    "[--render-dir DIR] [--max-shapes N] [--offset N]"
)

HELP_TEXT = (
    "pptx_inspect <file> [mode]; default overview. N = 5 | 2,5,8 | 2,5,7-9.\n"
    "--qa N: post-edit gate, #id text + full render + defect list\n"
    "--qa N --boxes: same, render relabelled with #id boxes\n"
    "--slide N: shapes, box, ph, type, fit, [!] OVERSET/HIDDEN\n"
    "--layouts: placeholders + type; static = master text\n"
    "--compare F: vs template F, ends [QA: PASS/FAIL]\n"
    "--validate: package integrity, does it open at all\n"
    "--grid: with --render only, tile slides into one image\n"
    "--text --media --render --render-dir --max-shapes --offset"
)


def picture_summary(shape: BaseShape) -> str:
    image = embedded_image(shape)
    if image is None:
        return ""
    parts: List[str] = []
    filename = getattr(image, "filename", None)
    if filename:
        parts.append(filename)
    size = getattr(image, "size", None)
    if isinstance(size, tuple) and len(size) == 2:
        parts.append(f"{size[0]}x{size[1]}px")
    content_type = getattr(image, "content_type", None)
    if content_type:
        parts.append(content_type)
    return "  ".join(parts)


def chart_summary(shape: BaseShape) -> str:
    chart = shape.chart
    parts: List[str] = []
    chart_type = getattr(chart, "chart_type", None)
    if chart_type is not None:
        type_name = getattr(chart_type, "name", str(chart_type))
        parts.append(type_name.lower())
    series_count = sum(1 for _ in chart.series)
    parts.append(f"series:{series_count}")
    title = ""
    if getattr(chart, "has_title", False):
        try:
            title = flatten_text(chart.chart_title.text_frame.text or "").strip()
        except (AttributeError, ValueError):
            title = ""
    if not title:
        title = flatten_text(chart_title_via_zip(shape)).strip()
    if title:
        parts.append(f'title:"{ellipsize(title, TEXT_PREVIEW_LIMIT)}"')
    return "  ".join(parts)


def chart_title_via_zip(shape: BaseShape) -> str:
    """Fallback: resolve the chart part via the slide's rels and parse the
    DrawingML title with the shared ooxml helper."""
    chart = getattr(shape, "chart", None)
    if chart is None:
        return ""
    part = getattr(chart, "part", None)
    partname = getattr(part, "partname", None)
    if not partname:
        return ""
    chart_path = str(partname).lstrip("/")
    pkg = getattr(part, "package", None)
    pkg_path = getattr(pkg, "_path", None) or getattr(pkg, "path", None)
    if not pkg_path:
        return ""
    try:
        zf = zipfile.ZipFile(pkg_path)
    except (zipfile.BadZipFile, OSError):
        return ""
    with zf:
        return ooxml.parse_chart_title(zf, chart_path) or ""


def table_summary(shape: BaseShape) -> Tuple[str, List[str]]:
    table = shape.table
    rows = list(table.rows)
    nrows = len(rows)
    ncols = len(rows[0].cells) if rows else 0
    summary = f"{nrows}x{ncols}"
    cell_lines: List[str] = []
    for r, row in enumerate(rows):
        for c, cell in enumerate(row.cells):
            text = flatten_text(cell.text or "").strip()
            if not text:
                continue
            cell_lines.append(
                f"  ({r + 1},{c + 1}) {ellipsize(text, TEXT_PREVIEW_LIMIT)}"
            )
    return summary, cell_lines


def slide_title(slide: Slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is None or not title_shape.has_text_frame:
        return ""
    return flatten_text(title_shape.text_frame.text or "").strip()


def slide_is_hidden(slide: Slide) -> bool:
    return slide.element.get("show") == "0"


def count_shapes_by_kind(shapes: Iterable[BaseShape]) -> dict:
    counts = {"text": 0, "pic": 0, "chart": 0, "table": 0, "other": 0}
    for shape in shapes:
        if shape.has_chart:
            counts["chart"] += 1
        elif shape.has_table:
            counts["table"] += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            counts["pic"] += 1
        elif _has_embedded_blip(getattr(shape, "_element", None)):
            counts["pic"] += 1  # picture placeholder / picture-filled shape
        elif shape.has_text_frame and (shape.text_frame.text or "").strip():
            counts["text"] += 1
        else:
            counts["other"] += 1
    return counts


def shape_is_hidden(shape: BaseShape) -> bool:
    """True when PowerPoint hides the shape from every render."""
    for el in shape._element.iter():
        if el.tag.endswith("}cNvPr"):
            return el.get("hidden") in ("1", "true")
    return False


def describe_shape(
    shape: BaseShape,
    *,
    ctx: Optional[SlideContext] = None,
    cover_candidates: Optional[List[CoverRect]] = None,
    all_boxes: Optional[List[CoverRect]] = None,
    layout_chain=None,
    indent: str = "",
) -> List[str]:
    kind = shape_kind(shape)
    box = format_box(shape)
    parts = [pad(f"#{shape.shape_id}", 6), pad(kind, 6), pad(box, 24)]
    ph = placeholder_type(shape)
    if ph:
        parts.append(pad(f"ph={ph}", 14))

    sub_lines: List[str] = []
    summary = ""
    if kind == "chart":
        summary = chart_summary(shape)
    elif kind == "pic":
        summary = picture_summary(shape)
    elif kind == "table":
        ts, cell_lines = table_summary(shape)
        summary = ts
        sub_lines.extend(cell_lines)
    elif kind == "group":
        nested = list(shape.shapes)
        summary = f"shapes:{len(nested)}"
        for inner in nested:
            # Cover detection doesn't translate into groups: children's
            # coordinates are group-local, and "siblings" inside the group
            # are typically arranged on purpose. Drop the candidate list.
            sub_lines.extend(
                describe_shape(inner, ctx=ctx, layout_chain=layout_chain, indent="    ")
            )
        if shape.has_text_frame:
            sub_lines.extend(text_frame_lines(shape, indent="  "))
    else:
        text_lines = text_frame_lines(shape, indent="  ")
        if text_lines:
            sub_lines.extend(text_lines)
        else:
            name = (shape.name or "").strip()
            if name:
                summary = f'name:"{ellipsize(name, 40)}"'

    if summary:
        parts.append(summary)
    # Vertical anchor (how text sits in its box): surfaced so floating text - a
    # short string anchored top in a much taller box - is diagnosable alongside
    # the boxed render. Explicit value when set; text boxes/auto shapes default
    # to top; placeholders inherit (left unstated).
    if kind != "group" and shape.has_text_frame:
        has_text = any(
            (p.text or "").strip() for p in shape.text_frame.paragraphs
        )
        if has_text:
            va = getattr(shape.text_frame, "vertical_anchor", None)
            va_name = getattr(va, "name", None)
            if va_name:
                parts.append(f"vanchor={va_name.lower()}")
            elif not ph:
                parts.append("vanchor=top")
    if shape_is_hidden(shape):
        parts.append("[!] HIDDEN (never renders; clear hidden=\"1\" to use it)")
    for marker in _shape_warning_markers(
        shape, ph, ctx, cover_candidates, all_boxes
    ):
        parts.append(marker)
    for token in _fit_tokens(shape, layout_chain):
        parts.append(token)

    head = indent + "  ".join(parts).rstrip()
    return [head] + [indent + line for line in sub_lines]


TITLE_PH_TYPES = frozenset({"title", "ctr_title", "center_title", "centertitle"})
BODY_PH_TYPES = frozenset({"body", "subtitle", "obj"})


def _theme_summary_line(file_path: str) -> Optional[str]:
    """One-line theme summary for the overview: name, bg/tx/accent colors.
    Fonts are reported separately by `_deck_fonts_line` so the agent does
    not conflate the theme's fallback typeface with the deck's actual one.
    Returns None if the package isn't readable."""
    try:
        zf = zipfile.ZipFile(file_path)
    except (zipfile.BadZipFile, OSError):
        return None
    with zf:
        pres_rels = ooxml.read_xml(
            zf, ooxml.rels_path_for("ppt/presentation.xml"))
        if pres_rels is None:
            return None
        master_path: Optional[str] = None
        for rel in pres_rels.findall("pr:Relationship", ooxml.NS):
            if rel.attrib.get("Type", "").endswith("/slideMaster"):
                master_path = ooxml.resolve_rel_target(
                    ooxml.rels_path_for("ppt/presentation.xml"),
                    rel.attrib["Target"],
                )
                break
        if not master_path:
            return None
        master_xml = ooxml.read_xml(zf, master_path)
        _, theme_xml = _read_theme_for_master(zf, master_path)
        if theme_xml is None:
            return None
        clr_map = _read_clr_map(master_xml)
        theme_colors = ooxml.theme_colors_by_name(theme_xml)
        theme_name = theme_xml.attrib.get("name") or "?"

        def _resolved(token: str) -> str:
            hx = _resolve_scheme_color(token, clr_map, theme_colors)
            return f"#{hx}" if hx else "-"

        accents = " ".join(_resolved(f"accent{i}") for i in range(1, 7))
        parts = [
            f"theme:{theme_name}",
            f"bg1:{_resolved('bg1')}",
            f"tx1:{_resolved('tx1')}",
            f"accents:{accents}",
        ]
        return "[" + " | ".join(parts) + "]"


def _format_font_key(
    key: Tuple[Optional[str], Optional[float], bool], varies: bool
) -> Optional[str]:
    typeface, size_pt, bold = key
    parts: List[str] = []
    if typeface:
        parts.append(typeface)
    if size_pt is not None:
        if float(size_pt).is_integer():
            parts.append(f"{int(size_pt)}pt")
        else:
            parts.append(f"{size_pt:.1f}pt")
    if bold:
        parts.append("bold")
    if not parts:
        return None
    rendered = " ".join(parts)
    return f"{rendered} (varies)" if varies else rendered


def _deck_fonts_line(prs: PresentationType, file_path: str) -> Optional[str]:
    """Resolve the dominant title and body placeholder typography across all
    layouts and report them alongside the theme's major/minor fallback.

    The theme's `major:`/`minor:` typefaces from <a:fontScheme> are only what
    runs outside placeholders inherit - many decks (e.g. the Dust template)
    declare Arial there but override every layout with Lexend. Reporting
    them as the deck's font misleads the agent into picking Arial for
    custom shapes. So we surface what the layouts actually resolve to, and
    keep the theme fallback only when it diverges from what the layouts use.
    """
    try:
        zf = zipfile.ZipFile(file_path)
    except (zipfile.BadZipFile, OSError):
        return None

    title_counts: Dict[Tuple[Optional[str], Optional[float], bool], int] = {}
    body_counts: Dict[Tuple[Optional[str], Optional[float], bool], int] = {}
    fallback_major: Optional[str] = None
    fallback_minor: Optional[str] = None

    with zf:
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                layout_path = _layout_part_path(layout)
                if not layout_path:
                    continue
                (
                    layout_xml,
                    master_xml,
                    theme_xml,
                    clr_map,
                    theme_colors,
                ) = _read_layout_chain(zf, layout_path)
                if fallback_major is None:
                    fallback_major = ooxml.theme_font(theme_xml, "major")
                if fallback_minor is None:
                    fallback_minor = ooxml.theme_font(theme_xml, "minor")
                if layout_xml is None:
                    continue
                for ph in layout.placeholders:
                    ph_name = placeholder_type(ph) or ""
                    if ph_name not in TITLE_PH_TYPES and ph_name not in BODY_PH_TYPES:
                        continue
                    pf = ph.placeholder_format
                    ph_idx = pf.idx if pf else None
                    defaults = resolve_placeholder_defaults(
                        layout_xml,
                        master_xml,
                        theme_xml,
                        clr_map,
                        theme_colors,
                        ph_idx=ph_idx,
                        ph_type=ph_name,
                    )
                    key = (
                        defaults.get("typeface"),
                        defaults.get("size_pt"),
                        bool(defaults.get("bold")),
                    )
                    bucket = title_counts if ph_name in TITLE_PH_TYPES else body_counts
                    bucket[key] = bucket.get(key, 0) + 1

    def _dominant(
        counts: Dict[Tuple[Optional[str], Optional[float], bool], int],
    ) -> Tuple[Optional[Tuple[Optional[str], Optional[float], bool]], bool]:
        if not counts:
            return None, False
        top_key = max(counts.items(), key=lambda kv: kv[1])[0]
        return top_key, len(counts) > 1

    title_key, title_varies = _dominant(title_counts)
    body_key, body_varies = _dominant(body_counts)

    title_str = _format_font_key(title_key, title_varies) if title_key else None
    body_str = _format_font_key(body_key, body_varies) if body_key else None

    title_face = title_key[0] if title_key else None
    body_face = body_key[0] if body_key else None
    title_matches_fallback = (
        fallback_major is not None and title_face == fallback_major
    )
    body_matches_fallback = (
        fallback_minor is not None and body_face == fallback_minor
    )
    show_fallback = bool(fallback_major or fallback_minor) and not (
        title_matches_fallback and body_matches_fallback
    )

    parts: List[str] = []
    if title_str:
        parts.append(f"title={title_str}")
    if body_str:
        parts.append(f"body={body_str}")
    if show_fallback:
        parts.append(
            f"theme-fallback={fallback_major or '-'}/{fallback_minor or '-'}"
        )
    if not parts:
        return None
    return "[fonts: " + " | ".join(parts) + "]"


def print_overview(prs: PresentationType, file_path: str) -> str:
    width = emu_to_inches(prs.slide_width) or 0.0
    height = emu_to_inches(prs.slide_height) or 0.0
    slide_count = len(prs.slides)
    layout_count = sum(len(m.slide_layouts) for m in prs.slide_masters)

    word_counts = [slide_word_count(slide) for slide in prs.slides]
    avg_words = sum(word_counts) // len(word_counts) if word_counts else 0
    max_words = max(word_counts) if word_counts else 0

    lines = [
        f"[Slides: {slide_count} | "
        f"size: {width:.1f}x{height:.1f}\" | "
        f"masters: {len(prs.slide_masters)} | "
        f"layouts: {layout_count} | "
        f"words/slide: avg={avg_words} max={max_words}]"
    ]
    theme_line = _theme_summary_line(file_path)
    if theme_line:
        lines.append(theme_line)
    fonts_line = _deck_fonts_line(prs, file_path)
    if fonts_line:
        lines.append(fonts_line)
    for idx, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout.name or "?"
        title = slide_title(slide)
        counts = count_shapes_by_kind(slide.shapes)
        words = word_counts[idx - 1]
        flags: List[str] = []
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            flags.append("notes")
        if slide_is_hidden(slide):
            flags.append("hidden")

        title_str = f'"{ellipsize(title, 40)}"' if title else ""
        counts_str = (
            f"shapes:{sum(counts.values())}  "
            f"text:{counts['text']}  pic:{counts['pic']}  "
            f"chart:{counts['chart']}  table:{counts['table']}  "
            f"words:{words}"
        )
        flags_str = f"  [{','.join(flags)}]" if flags else ""
        lines.append(
            f"  {pad(str(idx), 3)} {pad(ellipsize(layout, 24), 26)}"
            f"{pad(title_str, 42)}  {counts_str}{flags_str}"
        )
    return "\n".join(lines)


def find_slide(prs: PresentationType, idx: int) -> Slide:
    if idx < 1 or idx > len(prs.slides):
        raise ValueError(
            f"slide index out of range: {
                idx} (deck has {len(prs.slides)} slides)"
        )
    return prs.slides[idx - 1]


def print_slide(
    prs: PresentationType,
    file_path: str,
    idx: int,
    offset: int,
    max_shapes: int,
) -> str:
    slide = find_slide(prs, idx)
    shapes = list(slide.shapes)
    total = len(shapes)
    layout = slide.slide_layout.name or "?"
    title = slide_title(slide)
    title_str = f' | "{ellipsize(title, 60)}"' if title else ""
    flags: List[str] = []
    if slide_is_hidden(slide):
        flags.append("hidden")
    flags_str = f" | {','.join(flags)}" if flags else ""

    header = (
        f"[Slide {idx}/{len(prs.slides)} | layout: {layout}"
        f"{title_str}{flags_str} | shapes: {total}]"
    )
    lines = [header]

    ctx = SlideContext(
        width_emu=prs.slide_width or 0,
        height_emu=prs.slide_height or 0,
    )
    cover_candidates = _cover_candidates(shapes)
    all_boxes: List[CoverRect] = [
        (s.shape_id, s.left, s.top, s.width, s.height)
        for s in shapes
        if None not in (s.left, s.top, s.width, s.height)
    ]

    # Resolve the slide's layout chain once so fit estimates can use each
    # placeholder's inherited font size. Parsed elements outlive the zip.
    layout_chain = _resolve_layout_chain(file_path, slide)

    end = min(total, offset + max_shapes)
    for shape in shapes[offset:end]:
        lines.extend(
            describe_shape(
                shape,
                ctx=ctx,
                cover_candidates=cover_candidates,
                all_boxes=all_boxes,
                layout_chain=layout_chain,
            )
        )

    if end < total:
        lines.append("")
        lines.append(
            f"[Showing shapes {offset + 1}-{end} of {total}. "
            f"Use --offset {end} for the next page.]"
        )

    if slide.has_notes_slide:
        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            lines.append("")
            lines.append("[Notes]")
            for paragraph in notes.split("\n"):
                paragraph = flatten_text(paragraph).strip()
                if paragraph:
                    lines.append(
                        f"  {ellipsize(paragraph, TEXT_PREVIEW_LIMIT)}")
    return "\n".join(lines)


def print_slides(
    prs: PresentationType,
    file_path: str,
    slide_nos: List[int],
    offset: int,
    max_shapes: int,
) -> str:
    """Structural view of one or more slides (from a --slide N[,N,...] pattern),
    one block per slide, so a whole batch can be inspected before editing in a
    single call rather than one call each."""
    return "\n\n".join(
        print_slide(prs, file_path, idx, offset, max_shapes)
        for idx in slide_nos
    )


def _layout_part_path(layout) -> Optional[str]:
    """Layout XML path inside the .pptx zip, e.g. 'ppt/slideLayouts/slideLayout3.xml'."""
    partname = getattr(getattr(layout, "part", None), "partname", None)
    if partname is None:
        return None
    return str(partname).lstrip("/")


def _resolve_layout_chain(file_path: str, slide: Slide):
    """The slide's (layout, master, theme, clr_map, theme_colors) chain, used to
    resolve placeholder-inherited font sizes. Parsed elements outlive the zip.
    None when the package or its layout part can't be read."""
    try:
        zf: Optional[zipfile.ZipFile] = zipfile.ZipFile(file_path)
    except (zipfile.BadZipFile, OSError):
        return None
    with zf:
        layout_path = _layout_part_path(slide.slide_layout)
        if not layout_path:
            return None
        return _read_layout_chain(zf, layout_path)


def _text_extent_boxes(file_path: str, slide: Slide) -> Dict[int, Tuple[int, int, int, int]]:
    """Map shape_id -> the box (EMU) grown to wrap a text shape's rendered copy,
    for the shapes whose copy overflows the declared box. Feeds the --qa overlay
    so it wraps the actual text and catches overflow-into-neighbour overlaps.
    Empty when nothing overflows."""
    layout_chain = _resolve_layout_chain(file_path, slide)
    out: Dict[int, Tuple[int, int, int, int]] = {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not any((p.text or "").strip() for p in shape.text_frame.paragraphs):
            continue
        size_pt = _effective_font_size_pt(shape, layout_chain)
        if size_pt is None:
            continue
        box = _text_extent_box(shape, size_pt)
        if box is not None:
            out[shape.shape_id] = box
    return out


def _static_text_lines(shapes: Iterable[BaseShape], indent: str) -> List[str]:
    """Master/layout text shapes every inheriting slide renders."""
    lines: List[str] = []
    for shape in shapes:
        if shape.is_placeholder or not shape.has_text_frame:
            continue
        text = flatten_text(shape.text_frame.text).strip()
        if not text:
            continue
        lines.append(
            f"{indent}static  #{shape.shape_id}  {pad(format_box(shape), 24)}"
            f'  "{ellipsize(text, 60)}"'
        )
    return lines


def print_layouts(prs: PresentationType, file_path: str) -> str:
    lines = [f"[Masters: {len(prs.slide_masters)}]"]
    try:
        zf: Optional[zipfile.ZipFile] = zipfile.ZipFile(file_path)
    except (zipfile.BadZipFile, OSError):
        zf = None
    try:
        for mi, master in enumerate(prs.slide_masters, start=1):
            master_name = (master.name or "").strip() or f"master{mi}"
            lines.append("")
            lines.append(
                f"# Master {mi}: {master_name}  layouts: {
                    len(master.slide_layouts)}"
            )
            master_static = _static_text_lines(master.shapes, "  ")
            if master_static:
                lines.append(
                    "  [inherited text - renders on every slide using this "
                    "master, editable only on the master itself:]"
                )
                lines.extend(master_static)
            for layout in master.slide_layouts:
                placeholders = list(layout.placeholders)
                lines.append(
                    f"- {pad(layout.name or '?', 28)
                         } placeholders: {len(placeholders)}"
                )
                lines.extend(_static_text_lines(layout.shapes, "    "))

                layout_path = _layout_part_path(layout)
                layout_xml = master_xml = theme_xml = None
                clr_map: Dict[str, str] = {}
                theme_colors: Dict[str, str] = {}
                if zf is not None and layout_path:
                    (
                        layout_xml,
                        master_xml,
                        theme_xml,
                        clr_map,
                        theme_colors,
                    ) = _read_layout_chain(zf, layout_path)

                for ph in placeholders:
                    ph_name = placeholder_type(ph) or "unknown"
                    pf = ph.placeholder_format
                    idx = pf.idx if pf else None
                    idx_str = str(idx) if idx is not None else "?"
                    box = format_box(ph)
                    head = f"    [{idx_str}] {pad(ph_name, 14)} {pad(box, 24)}"
                    if shape_is_hidden(ph):
                        head = f'{head}  [!] HIDDEN (never renders; clear hidden="1" to use it)'
                    if layout_xml is not None:
                        defaults = resolve_placeholder_defaults(
                            layout_xml,
                            master_xml,
                            theme_xml,
                            clr_map,
                            theme_colors,
                            ph_idx=idx,
                            ph_type=ph_name,
                        )
                        defaults_str = format_placeholder_defaults(defaults)
                        if defaults_str:
                            head = f"{head}  {defaults_str}"
                    lines.append(head.rstrip())
    finally:
        if zf is not None:
            zf.close()
    return "\n".join(lines)


# Text that looks like un-filled template scaffolding: bracketed/angled
# prompts, lorem/placeholder fillers. Repeated-across-slides copy (e.g.
# "Subject title", "Summary") is detected separately in print_text.
def print_text(prs: PresentationType, slide_idx: Optional[int] = None) -> str:
    # Pass 1: find distinctive copy repeated across the deck - template
    # scaffolding the author forgot to replace ("Subject title", "Summary",
    # "Title of the slide"). Count total occurrences (catches repeats within a
    # single slide too); require length >= 8 so content words ("Dust") don't trip.
    repeat_counts: Dict[str, int] = {}
    for slide in prs.slides:
        for shape in slide.shapes:
            for raw in _shape_text_iter(shape):
                t = flatten_text(raw).strip()
                if len(t) >= 8:
                    repeat_counts[t] = repeat_counts.get(t, 0) + 1
    repeated = sorted(
        (t for t, c in repeat_counts.items() if c >= 3),
        key=lambda t: -repeat_counts[t],
    )

    blocks: List[str] = []
    total_chars = 0
    for idx, slide in enumerate(prs.slides, start=1):
        if slide_idx is not None and idx != slide_idx:
            continue
        slide_lines: List[str] = []
        for shape in slide.shapes:
            slide_lines.extend(_collect_text(shape))
        if slide.has_notes_slide:
            notes = flatten_text(
                slide.notes_slide.notes_text_frame.text or ""
            ).strip()
            if notes:
                slide_lines.append(
                    f"  [notes] {ellipsize(notes, TEXT_PREVIEW_LIMIT)}")
        if slide_lines:
            title = slide_title(slide)
            header = f"# Slide {idx} (words:{slide_word_count(slide)})"
            if title:
                header += f': "{ellipsize(title, 60)}"'
            blocks.append(header)
            blocks.extend(slide_lines)
            blocks.append("")
            total_chars += sum(len(line) for line in slide_lines)
    if not blocks:
        return "[No text in deck]"
    if blocks and blocks[-1] == "":
        blocks.pop()
    scope = (f"slide {slide_idx}" if slide_idx is not None
             else f"{len(prs.slides)} slides")
    head = f"[Text: {total_chars} chars | {scope} | each line is tagged with " \
           "its shape #id - match against the --qa box labels]"
    if repeated:
        shown = ", ".join(
            f'"{ellipsize(t, 30)}" x{repeat_counts[t]}' for t in repeated[:6]
        )
        head += f"\n[i] copy repeated 3+ times: {shown}"
    return head + "\n\n" + "\n".join(blocks)


def _collect_text(shape: BaseShape, indent: str = "  ") -> List[str]:
    lines: List[str] = []
    kind = shape_kind(shape)
    if kind == "group":
        for inner in shape.shapes:
            lines.extend(_collect_text(inner, indent))
        return lines
    sid = shape.shape_id
    if kind == "table":
        _, cell_lines = table_summary(shape)
        lines.extend(f"{indent}#{sid} {cl.strip()}" for cl in cell_lines)
        return lines
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = flatten_text(paragraph.text or "").strip()
            if not text:
                continue
            level = paragraph.level or 0
            mark = " [leftover?]" if _is_leftover_suspect(text) else ""
            lines.append(
                f"{indent}#{sid} p{level}: "
                f"{ellipsize(text, TEXT_PREVIEW_LIMIT)}{mark}"
            )
    return lines


# Findings the deck-level `shapes:` line already reports, baselined against the
# template. Naming them here keeps the roll-up to the checks that have no
# deck-level equivalent (legibility, artwork, clipping, collisions, overset).
_BASELINED_PHRASES = (
    "zero-size box",
    "extends past slide edge",
    "image distorted",
    "stacked with shape",
)


def _deck_slide_defects(file_path: str) -> Optional[List[Tuple[int, str]]]:
    """Every `[!]` the per-slide QA gate would raise, across the deck, as
    (slide_no, line). None when the deck could not be rendered.

    Runs the SAME code path as `--qa` rather than a deck-level approximation, so
    the final audit can never pass a slide the per-slide gate would fail. It used
    to: collisions, clipping and the shape markers were computed only in `--qa`,
    so a deck whose title slide was three headlines printed on top of each other
    still ended in [QA: PASS] if nobody ran `--qa` on slide 1.

    The four structural kinds are left out: `shapes:` above already reports them
    against the template's own, and a deck that cloned an exemplar carrying five
    zero-size boxes should not be billed for them twice - once baselined and
    once not.

    Rendered a page at a time so the soffice PDF the QA pass already wrote is
    reused: after QA this costs a rasterization per slide, not a reconversion."""
    try:
        prs = Presentation(file_path)
        out: List[Tuple[int, str]] = []
        for slide_idx in range(1, len(prs.slides) + 1):
            _image, digest = _annotate_slide(file_path, prs, slide_idx)
            for line in digest:
                stripped = line.strip()
                if not stripped.startswith("[!]") or line.startswith("[!]"):
                    continue
                if any(phrase in stripped for phrase in _BASELINED_PHRASES):
                    continue
                out.append((slide_idx, stripped))
        return out
    except (ValueError, OSError, IndexError):
        return None


def print_compare(file_path: str, source_path: str) -> str:
    """Compare the edited deck (file_path) against its source/template
    (source_path) and gate the result. Surfaces the regressions that mean the
    deck no longer respects the template: orphaned slide parts (hand-deleted
    instead of via pptx_slides), embedded media or fonts dropped, the template's
    imagery stripped, slides collapsed onto one catch-all layout, or density
    blown past the template's ceiling. Ends with a [QA: PASS/FAIL] verdict -
    do not deliver until it reads PASS."""
    out_media = _package_names(file_path, "ppt/media/")
    src_media = _package_names(source_path, "ppt/media/")
    out_fonts = _package_names(file_path, "ppt/fonts/")
    src_fonts = _package_names(source_path, "ppt/fonts/")
    if None in (out_media, src_media, out_fonts, src_fonts):
        return "[Compare: could not read one of the files as a .pptx package]"

    def delta(out_n: int, src_n: int) -> str:
        d = out_n - src_n
        return f"{d:+d}" if d else "0"

    blockers = 0

    # Parse each deck independently so one failing to parse doesn't silently
    # blind the whole audit. python-pptx reads only the slide list; the zip may
    # hold more parts.
    def _fid(path: str) -> Optional[DeckFidelity]:
        try:
            return _deck_fidelity(Presentation(path))
        except Exception:  # noqa: BLE001 - degrade visibly rather than crash
            return None

    out_fid = _fid(file_path)
    src_fid = _fid(source_path)

    out_zip = _count_slides(file_path)
    src_zip = _count_slides(source_path)
    out_listed = out_fid.total if out_fid else (_listed_slide_count(file_path) or out_zip)
    src_listed = src_fid.total if src_fid else (_listed_slide_count(source_path) or src_zip)

    lines = [
        f"[Compare: {os.path.basename(file_path)}  vs source "
        f"{os.path.basename(source_path)}]",
        f"  slides:  {pad(str(out_listed), 4)} (source {src_listed})  "
        f"{delta(out_listed, src_listed)}",
    ]

    # Orphaned slide parts: in the package but not in the slide list. This is
    # the signature of deleting slides by hand (removing sldIds) instead of
    # `pptx_slides --delete`, which would drop the parts and their media too.
    orphans = out_zip - out_listed
    if orphans > 0:
        blockers += 1
        lines.append(
            f"  orphans: [!] {orphans} slide part(s) in the package but not in the "
            "slide list (orphaned)"
        )

    # A deck we cannot parse cannot be audited - say so loudly and block, so a
    # parse failure can never masquerade as a clean PASS.
    if out_fid is None or src_fid is None:
        blockers += 1
        which = "the edited deck" if out_fid is None else "the template"
        lines.append(
            f"  audit:   [!] degraded - could not parse {which} with python-pptx; "
            "imagery / bare-canvas / density not checked"
        )

    dropped_media = sorted(src_media - out_media)
    lines.append(
        f"  media:   {pad(str(len(out_media)), 4)} (source "
        f"{len(src_media)})  {delta(len(out_media), len(src_media))}"
    )
    if dropped_media:
        lines.append(
            f"    {len(dropped_media)} media in source not in output:"
        )
        for name in dropped_media:
            lines.append(f"      - {name[len('ppt/media/'):]}")

    dropped_fonts = sorted(src_fonts - out_fonts)
    lines.append(
        f"  fonts:   {pad(str(len(out_fonts)), 4)} (source "
        f"{len(src_fonts)})  {delta(len(out_fonts), len(src_fonts))}"
    )
    if dropped_fonts:
        blockers += 1
        lines.append(
            f"    [!] {len(dropped_fonts)} embedded font part(s) in source not in "
            "output:"
        )
        for name in dropped_fonts:
            lines.append(f"      - {name[len('ppt/fonts/'):]}")

    if out_fid and src_fid and out_fid.total and src_fid.total:
        audit_content = out_fid.total > MIN_SLIDES_FOR_CONTENT_AUDIT

        # Imagery: did the rendered slides keep the template's pictures/charts/
        # tables, or get rebuilt as text on the background?
        src_rate = src_fid.imagery_slides / src_fid.total
        out_rate = out_fid.imagery_slides / out_fid.total
        lines.append(
            f"  imagery: pics/charts/tables on {out_fid.imagery_slides}/"
            f"{out_fid.total} slides ({out_fid.imagery_objs} objects; template "
            f"{src_fid.imagery_slides}/{src_fid.total})"
        )
        if audit_content and src_rate >= IMAGERY_RICH_TEMPLATE and out_rate < src_rate * IMAGERY_STRIP_RATIO:
            blockers += 1
            lines.append(
                f"    [!] imagery below template - {out_fid.imagery_slides}/"
                f"{out_fid.total} output slides carry images vs "
                f"{src_fid.imagery_slides}/{src_fid.total} in the template. "
                "Clone the template's own image-bearing slides for the points "
                "that need one; do not paste the same picture onto every text "
                "slide to raise the count."
            )

        # Bare canvas: slides hand-drawn with no template placeholders and no
        # imagery - the signature of rebuilding on a blank layout. Judged
        # relative to the template, since some templates ship shape-built slides.
        out_bare = out_fid.bare_slides / out_fid.total
        src_bare = src_fid.bare_slides / src_fid.total
        lines.append(
            f"  layouts: {len(out_fid.layout_counts)} used; "
            f"{out_fid.bare_slides}/{out_fid.total} slides on an empty canvas "
            f"(template {src_fid.bare_slides}/{src_fid.total})"
        )
        if audit_content and out_bare >= BARE_RATE and out_bare > src_bare + BARE_MARGIN:
            blockers += 1
            lines.append(
                f"    [!] {out_fid.bare_slides}/{out_fid.total} slides have no "
                f"placeholders and no imagery (template {src_fid.bare_slides}/"
                f"{src_fid.total})"
            )

        # Density: the template's max words/slide is a hard ceiling. A ceiling of
        # zero means the template ships layouts with no copy on them, so it sets
        # no budget and every filled slide would "exceed" it.
        ceiling = max(src_fid.word_counts) if src_fid.word_counts else 0
        if out_fid.word_counts and ceiling:
            out_max = max(out_fid.word_counts)
            over = [
                (i + 1, w)
                for i, w in enumerate(out_fid.word_counts)
                if w > ceiling * DENSITY_TOLERANCE
            ]
            lines.append(
                f"  density: max {out_max} words/slide "
                f"(template ceiling {ceiling})"
            )
            if over:
                blockers += 1
                listing = ", ".join(f"{i}({w})" for i, w in over)
                lines.append(
                    f"    [!] {len(over)} slide(s) over template ceiling of "
                    f"{ceiling} words: {listing}"
                )
        elif out_fid.word_counts:
            lines.append(
                f"  density: max {max(out_fid.word_counts)} words/slide "
                "(template sets no ceiling - its slides carry no copy)"
            )

        # Per-slide structural issues, rolled up and BASELINED against the
        # template's own. Reporting these without gating on them was how a deck
        # with five boxes off the slide edge and four stretched photos still
        # read [QA: PASS]; baselining is what lets them gate without blaming the
        # model for faults the template shipped with.
        structural = _deck_structural_audit(file_path, source_path)
        if structural is not None:
            summary = " ".join(
                f"{STRUCTURAL_LABELS[kind]}:"
                f"{sum(1 for _, _, inh in structural[kind] if not inh)}"
                f"(+{sum(1 for _, _, inh in structural[kind] if inh)} from template)"
                for kind in STRUCTURAL_KINDS
            )
            lines.append(f"  shapes:  {summary}")
            for kind in STRUCTURAL_KINDS:
                own = [(n, sid) for n, sid, inh in structural[kind] if not inh]
                if not own:
                    continue
                blockers += len(own)
                where = ", ".join(f"slide {n} #{sid}" for n, sid in own[:8])
                more = (
                    f" and {len(own) - 8} more" if len(own) > 8 else ""
                )
                lines.append(
                    f"    [!] {len(own)} {STRUCTURAL_LABELS[kind]} shape(s) you "
                    f"added or moved: {where}{more}."
                )

    # Content-slot fidelity: text written into the template's interior spacer
    # paragraphs (the wrong-slot fill). Deterministic from the paragraph
    # skeletons - no estimation - so it gates as a blocker.
    slot_findings = _slot_audit(file_path, source_path)
    if slot_findings:
        lines.append(
            f"  slots:   [!] {len(slot_findings)} text box(es) filled the "
            "template's spacer paragraphs"
        )
        for slide_no, sid, spacers, content in slot_findings:
            blockers += 1
            filled = ", ".join(f"p[{i}]" for i in sorted(spacers))
            slots = ", ".join(f"p[{i}]" for i in sorted(content))
            lines.append(
                f"    [!] slide {slide_no} #{sid}: filled spacer {filled}; "
                f"template content slots {slots}"
            )

    # Legibility, measured on the render: text the same colour family as what
    # ends up behind it. python-pptx cannot see this - a text box added outside
    # a placeholder takes the presentation default, which is near-black whatever
    # the slide looks like - and it is the most common way an edited deck ships
    # slides nobody can read.
    if out_fid is not None:
        defects = _deck_slide_defects(file_path)
        if defects is None:
            lines.append(
                "  defects: [i] could not render the deck; the per-slide checks "
                "(legibility, collisions, clipping) did not run"
            )
        elif defects:
            lines.append(
                f"  defects: [!] {len(defects)} per-slide defect(s) - the same "
                "ones --qa raises"
            )
            blockers += len(defects)
            for slide_no, line in defects[:LEFTOVER_LISTED]:
                lines.append(f"    [!] slide {slide_no}: {line[len('[!] '):]}")
            if len(defects) > LEFTOVER_LISTED:
                lines.append(
                    f"    [!] ... and {len(defects) - LEFTOVER_LISTED} more; run "
                    "--qa on each slide for its own list"
                )
        else:
            lines.append("  defects: every slide passes the per-slide checks")

    # Template filler still in the deck. No template needed and no judgment
    # call: lorem or a bracketed prompt in a delivered deck is always a defect.
    filler = _filler_audit(file_path)
    if filler:
        lines.append(
            f"  filler:  [!] {len(filler)} shape(s) still hold template filler"
        )
        for slide_no, sid, text in filler[:LEFTOVER_LISTED]:
            blockers += 1
            lines.append(
                f"    [!] slide {slide_no} #{sid}: {ellipsize(text, 60)!r}"
            )
        if len(filler) > LEFTOVER_LISTED:
            blockers += len(filler) - LEFTOVER_LISTED
            lines.append(
                f"    [!] ... and {len(filler) - LEFTOVER_LISTED} more"
            )

    # Blank rows left at the foot of a cloned table.
    blank_rows = _empty_table_rows_audit(file_path)
    if blank_rows:
        lines.append(
            f"  tables:  [!] {len(blank_rows)} table(s) end in empty rows"
        )
        for slide_no, shape_id, trailing, total in blank_rows[:LEFTOVER_LISTED]:
            blockers += 1
            lines.append(
                f"    [!] slide {slide_no} #{shape_id}: the last {trailing} of "
                f"{total} rows are empty and still draw their fill. Delete the "
                "surplus rows instead of blanking them."
            )

    # A line break at the start of a paragraph.
    stranded = _leading_break_audit(file_path)
    if stranded:
        lines.append(
            f"  breaks:  [!] {len(stranded)} paragraph(s) start with a line break"
        )
        for slide_no, shape_id, index, text in stranded[:LEFTOVER_LISTED]:
            blockers += 1
            lines.append(
                f"    [!] slide {slide_no} #{shape_id} p[{index}]: "
                f"{ellipsize(text, 46)!r} - starts on the second line of its "
                "own paragraph, and a bulleted one strands its bullet above it. "
                "Drop the leading break."
            )
        if len(stranded) > LEFTOVER_LISTED:
            blockers += len(stranded) - LEFTOVER_LISTED
            lines.append(
                f"    [!] ... and {len(stranded) - LEFTOVER_LISTED} more"
            )

    # The same sentence on several slides. Padding to satisfy a gate looks
    # exactly like this, and so does a template paragraph nobody rewrote.
    repeated = _repeated_text_audit(file_path)
    if repeated:
        lines.append(
            f"  repeats: [!] {len(repeated)} block(s) of copy appear more than "
            "once"
        )
        for text, slides in repeated[:LEFTOVER_LISTED]:
            blockers += 1
            where = ",".join(str(n) for n in slides)
            lines.append(
                f"    [!] on slide(s) {where}: {ellipsize(text, 60)!r}"
            )
        if len(repeated) > LEFTOVER_LISTED:
            blockers += len(repeated) - LEFTOVER_LISTED
            lines.append(
                f"    [!] ... and {len(repeated) - LEFTOVER_LISTED} more"
            )

    # One picture doing duty as the imagery of many slides.
    padded = _repeated_image_audit(file_path, source_path)
    if padded:
        for _key, slides in padded[:LEFTOVER_LISTED]:
            blockers += 1
            where = ",".join(str(n) for n in slides)
            lines.append(
                f"  padding: [!] one picture is the content image of slides "
                f"{where}. Give a slide the template's own image for what it "
                "says, or leave it text-only."
            )

    # Cloned slides that came out mostly empty canvas. The shape-retention
    # advisory below counts shapes; this measures the hole, which is what the
    # reader sees - dropping one full-bleed photo keeps most of the shapes and
    # loses most of the slide.
    holes = _hole_audit(file_path, source_path)
    if holes:
        lines.append(
            f"  canvas:  [!] {len(holes)} cloned slide(s) left mostly empty"
        )
        for out_no, src_no, cov, src_cov in holes[:LEFTOVER_LISTED]:
            blockers += 1
            lines.append(
                f"    [!] slide {out_no} covers {cov:.0%} of the canvas; its "
                f"exemplar (template slide {src_no}) covers {src_cov:.0%}. "
                "Put the exemplar's content back, or clone a sparser one."
            )

    # Exemplar copy still on a cloned slide: the template's scaffolding words
    # ("Pilot", "HOW", "01".."06", a stage label) shipped as if they were
    # content. Advisory: revising the user's own deck legitimately keeps most of
    # its copy, so the model judges each one.
    leftovers = _leftover_copy_audit(file_path, source_path)
    untouched = untouched_slides(leftovers, file_path)
    if untouched:
        lines.append(
            f"  cloned:  [!] {len(untouched)} slide(s) are still the template's "
            "slide, with their copy unchanged"
        )
        for slide_no, kept, total in untouched[:LEFTOVER_LISTED]:
            blockers += 1
            lines.append(
                f"    [!] slide {slide_no}: {kept} of {total} text shapes still "
                "hold the exemplar's copy. Write this slide, or delete it."
            )
    if leftovers:
        lines.append(
            f"  leftover: [i] {len(leftovers)} shape(s) still carry the "
            "template's own copy - replace or delete the ones you cloned as a "
            "layout, keep the ones you were asked to leave alone"
        )
        for slide_no, sid, text in leftovers[:LEFTOVER_LISTED]:
            lines.append(f"    [i] slide {slide_no} #{sid}: {ellipsize(text, 60)!r}")
        if len(leftovers) > LEFTOVER_LISTED:
            lines.append(
                f"    [i] ... and {len(leftovers) - LEFTOVER_LISTED} more"
            )

    # Shape retention: a cloned slide that dropped most of its exemplar's shapes
    # has been gutted rather than adapted. Advisory ([i]), not a blocker -
    # trimming surplus is expected and heavy adaptation is sometimes right.
    drops = _drop_audit(file_path, source_path)
    if drops:
        lines.append(
            f"  reuse:   [i] {len(drops)} slide(s) kept "
            f"<{int(SHAPE_RETENTION_FLOOR * 100)}% of their exemplar's shapes"
        )
        for out_no, src_no, kept, dropped in drops:
            lines.append(
                f"    [i] slide {out_no} (from template slide {src_no}): "
                f"kept {kept}, dropped {dropped} exemplar shape(s)"
            )

    # Package integrity, baselined against the template so its own faults are
    # not reported as yours. Fidelity says nothing about whether the file opens.
    invalid, inherited = pptx_validate.check_against(file_path, source_path)
    if invalid or inherited:
        lines.append("")
        lines.append("Package:")
        lines.extend(f"  [!] {problem}" for problem in invalid)
        if inherited:
            lines.append(f"  [i] {inherited} problem(s) already in the template")
        blockers += len(invalid)

    lines.append("")
    if blockers:
        lines.append(
            f"[QA: FAIL - {blockers} blocker{'s' if blockers != 1 else ''} "
            "([!] above)]"
        )
    else:
        lines.append("[QA: PASS]")

    return "\n".join(lines)


def print_media(path: str) -> str:
    try:
        zf = zipfile.ZipFile(path)
    except zipfile.BadZipFile:
        return "[Not a valid zip / .pptx package]"
    media_entries: List[Tuple[str, int]] = []
    with zf:
        for info in zf.infolist():
            if info.filename.startswith("ppt/media/"):
                media_entries.append((info.filename, info.file_size))
    if not media_entries:
        return "[No embedded media]"
    media_entries.sort()
    lines = [f"[Media: {len(media_entries)}]"]
    for name, size in media_entries:
        short = name[len("ppt/media/"):]
        lines.append(f"- {pad(short, 32)} {format_size(size)}")
    return "\n".join(lines)


# pptx QA / preview renders are published onto the conversation/pod mount
# (where the model can open them with files__cat) under this dot-prefixed
# subdir; the shared publish machinery lives in render_publish.
_VIEW_SUBDIR = ".pptx_render"


def _collect_tokens(shape: BaseShape) -> set:
    """Comparison tokens of a shape's text (group child and table cell text
    folded in), used to attribute a rendered PDF word back to its source shape.

    Tables count: without them a table's rendered words fall through to the
    nearest-box fallback and get charged to whatever text shape sits closest,
    which mislabels the finding and drags that shape's contrast median with it."""
    tokens: set = set()
    for text in _shape_text_iter(shape):
        for word in text.split():
            tok = pdf_text.norm_token(word)
            if tok:
                tokens.add(tok)
    return tokens


def _shape_text_label(shape: BaseShape) -> str:
    """A short readable label for a text shape (first non-empty child for a
    group), used in the collision digest's `#id "text"` references."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            label = _shape_text_label(child)
            if label:
                return label
        return ""
    if getattr(shape, "has_text_frame", False):
        return ellipsize(flatten_text(shape.text_frame.text or "").strip(), 30)
    for text in _shape_text_iter(shape):
        stripped = flatten_text(text).strip()
        if stripped:
            return ellipsize(stripped, 30)
    return ""


def _autofit_off(shape: BaseShape) -> bool:
    """True when the shape explicitly won't auto-shrink its text (noAutofit) -
    surfaced as a fix hint on the overflowing shape."""
    return (
        getattr(shape, "has_text_frame", False)
        and shape.text_frame.auto_size == MSO_AUTO_SIZE.NONE
    )


def _text_shape_entries(slide: Slide):
    """Per top-level text-bearing shape: (id, declared_box_emu, token_set, label,
    autofit_off, declared_word_count). Feeds the cross-shape overprint detector
    (id, box, tokens), the digest wording (label, autofit) and the clipped-text
    check (declared_word_count)."""
    out = []
    for sh in slide.shapes:
        if None in (sh.left, sh.top, sh.width, sh.height):
            continue
        tokens = _collect_tokens(sh)
        if not tokens:
            continue
        box = (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)
        declared = sum(len(text.split()) for text in _shape_text_iter(sh))
        out.append((
            sh.shape_id, box, tokens, _shape_text_label(sh), _autofit_off(sh),
            declared,
        ))
    return out


_OVERFLOW_EDGE_PHRASE = {
    "below": "below its box",
    "above": "above its box",
    "right": "past its right edge",
    "left": "past its left edge",
}


def _slide_findings_lines(
    slide: Slide,
    idx: int,
    findings: dict,
    words: Optional[List[pdf_text.WordBox]],
    contrasts: Optional[List[pptx_contrast.ShapeContrast]] = None,
    shape_blockers: Optional[List[Tuple[int, str]]] = None,
    slide_blockers: Optional[List[str]] = None,
) -> List[str]:
    """Format one slide's QA findings as a tiered, consequence-first block.

    Read straight off the renderer's word positions (`words` from the soffice
    PDF), in three tiers:
      [!] fix before delivery - two DIFFERENT shapes' strongly-attributed words
          overprint (`pdf_text.cross_shape_overprints`).
      [w] review - a shape's OWN words land outside its OWN box
          (`pdf_text.self_overflows`): the box doesn't contain its text. Usually
          a real defect, sometimes decorative overflow into empty space, so the
          model decides - it does NOT gate delivery.
      [i] FYI - an unconfirmed geometric spill, a box overlap the render couldn't
          read (text exported as curves), or a drifted decorative marker.
    Reading where the glyphs actually landed - instead of estimating wrap from
    char counts or arguing about box geometry - is what keeps false positives
    down: same-run overlaps (superscripts, footnote markers, emoji) and box
    overlaps in empty space never qualify as `[!]`. A clean slide says so."""
    entries = _text_shape_entries(slide)
    label = {e[0]: e[3] for e in entries}
    autofit_off = {e[0]: e[4] for e in entries}

    def q(sid: int) -> str:
        lbl = label.get(sid)
        return f'#{sid} "{lbl}"' if lbl else f"#{sid}"

    severe: List[str] = []
    review: List[str] = []
    advisory: List[str] = []
    n_severe = 0
    n_review = 0

    for sid, marker in shape_blockers or []:
        n_severe += 1
        severe.append(f"  [!] {q(sid)} - {marker}.")

    for marker in slide_blockers or []:
        n_severe += 1
        severe.append(f"  [!] {marker}.")

    contrast_severe, contrast_review, n_contrast = pptx_contrast.contrast_lines(
        contrasts or [], q
    )
    severe.extend(contrast_severe)
    review.extend(contrast_review)
    n_severe += n_contrast
    n_review += len(contrast_review)

    if words:
        shapes = [(e[0], e[1], e[2]) for e in entries]
        confirmed_pairs = set()
        confirmed_sids = set()
        for c in pdf_text.cross_shape_overprints(words, shapes):
            n_severe += 1
            over, hit = c["over"], c["hit"]
            confirmed_pairs.add(frozenset((over, hit)))
            confirmed_sids.update((over, hit))
            if c["symmetric"]:
                severe.append(
                    f"  [!] text-on-text - {q(over)} and {q(hit)} overprint "
                    f"(rendered \"{c['word_over']}\" lands on \"{c['word_hit']}\")."
                )
                severe.append(
                    "      Possible fixes: separate the two boxes, or "
                    "shrink/reposition one."
                )
            else:
                shrink = (
                    ", enable shrink-to-fit (it's set to no-autofit)"
                    if autofit_off.get(over)
                    else ""
                )
                severe.append(
                    f"  [!] text-on-text - {q(over)} overflows its box onto "
                    f"{q(hit)} (rendered \"{c['word_over']}\" sits on "
                    f"\"{c['word_hit']}\")."
                )
                severe.append(
                    f"      Possible fixes: shorten #{over}, lower its font "
                    f"size{shrink}, or move #{hit} clear."
                )
        # A geometric spill the render did NOT confirm is not necessarily clean.
        # Substitute fonts (the deck's real face is often absent at render time)
        # reflow the text so it never overflows on the rendered page, hiding a
        # real overprint the declared geometry still predicts. Surface such a
        # spill as a review item rather than letting the render silently clear
        # it. Peer overlaps are excluded - those are usually designed overlays,
        # exactly what reading the render is trusted to filter out.
        for f in findings.get("overlaps", []):
            if not f["text_on_text"] or f["kind"] != "spill":
                continue
            if frozenset((f["over"], f["hit"])) in confirmed_pairs:
                continue
            advisory.append(
                f"  [i] {q(f['over'])} overflows its box onto {q(f['hit'])} by "
                f"{round(f['coverage'] * 100)}% at declared geometry, but the "
                "render didn't confirm text-on-text - if the deck's fonts were "
                "substituted the render can understate the overlap; verify."
            )
        # A shape whose own rendered words fall outside its own box: the box
        # doesn't contain its text, but the overflow lands in empty space (a
        # neighbour hit is a [!] collision above, so those shapes are skipped).
        for f in pdf_text.clipped_shapes(
            words, shapes, {e[0]: e[5] for e in entries}, page_has_text=True
        ):
            n_severe += 1
            severe.append(
                f"  [!] {q(f['sid'])} - text clipped: {f['rendered']} of "
                f"{f['declared']} words rendered. The box cuts its copy off "
                "instead of overflowing, so nothing looks wrong in the "
                "geometry. Shorten the copy or grow the box."
            )

        for f in pdf_text.self_overflows(words, shapes):
            if f["sid"] in confirmed_sids:
                continue
            n_review += 1
            where = _OVERFLOW_EDGE_PHRASE[f["edge"]]
            review.append(
                f"  [w] {q(f['sid'])} - text runs ~{f['over_in']:.2f}in {where} "
                f"(rendered \"{f['word']}\" lands outside it)."
            )
            review.append(
                "      The box doesn't contain its text: resize it, or shorten / "
                "shrink-to-fit the copy. Fine if it's decorative overflow into "
                "empty space."
            )
    else:
        # No extractable rendered text (e.g. text exported as curves): can't
        # confirm against glyphs, so surface the box-overlap candidates for a
        # manual look rather than confirming or silently clearing them.
        for f in findings.get("overlaps", []):
            if not f["text_on_text"]:
                continue
            advisory.append(
                f"  [i] {q(f['over'])} and {q(f['hit'])} boxes overlap by "
                f"{round(f['coverage'] * 100)}% - couldn't read the rendered "
                "text to confirm; check the render."
            )

    for sid, off in findings.get("markers", []):
        advisory.append(
            f"  [i] {q(sid)} marker has no text row within {off:.2f}in - "
            "check its alignment in the render."
        )

    if severe:
        header = (
            f"[!] slide {idx} - {n_severe} defect"
            f"{'s' if n_severe != 1 else ''} to fix before delivery:"
        )
        return [header] + severe + review + advisory
    if review:
        header = (
            f"[w] slide {idx} - {n_review} item"
            f"{'s' if n_review != 1 else ''} to review:"
        )
        return [header] + review + advisory
    if advisory:
        return [
            f"[i] slide {idx} - {len(advisory)} item"
            f"{'s' if len(advisory) != 1 else ''} to review:"
        ] + advisory
    return [
        f"[slide {idx}: mechanical checks clean - the render is still the gate]"
    ]


def _slide_contrasts(
    file_path: str,
    prs: PresentationType,
    slide,
    slide_idx: int,
    image_path: Path,
    words,
) -> List[pptx_contrast.ShapeContrast]:
    entries = _text_shape_entries(slide)
    return pptx_contrast.shape_contrasts(
        image_path,
        words,
        [(e[0], e[1], e[2]) for e in entries],
        prs.slide_width or 0,
        prs.slide_height or 0,
    )


def _slide_shape_blockers(
    file_path: str, prs: PresentationType, slide, slide_idx: int
) -> List[Tuple[int, str]]:
    """The `[!]` markers the --slide view prints for each shape, as (id, text).

    The QA gate used to compute collisions only, so a box off the slide edge, a
    stretched photo, a collapsed box and an unfilled placeholder all rendered a
    clean "no collisions" while --slide flagged them - the model's post-edit
    check could not see the defects that most often ship."""
    ctx = SlideContext(
        width_emu=prs.slide_width or 0, height_emu=prs.slide_height or 0
    )
    shapes = list(slide.shapes)
    cover_candidates = _cover_candidates(shapes)
    all_boxes: List[CoverRect] = [
        (s.shape_id, s.left, s.top, s.width, s.height)
        for s in shapes
        if None not in (s.left, s.top, s.width, s.height)
    ]
    layout_chain = _resolve_layout_chain(file_path, slide)
    out: List[Tuple[int, str]] = []
    for shape in shapes:
        text = flatten_text(" ".join(_shape_text_iter(shape))).strip()
        if text and _is_leftover_suspect(text):
            out.append(
                (shape.shape_id, f"still holds template filler: {ellipsize(text, 50)!r}")
            )
        markers = _shape_warning_markers(
            shape,
            placeholder_type(shape),
            ctx,
            cover_candidates=cover_candidates,
            all_boxes=all_boxes,
        )
        markers.extend(_fit_tokens(shape, layout_chain))
        for marker in markers:
            if marker.startswith("[!]"):
                out.append((shape.shape_id, marker[len("[!] "):]))
    return out


def _annotate_slide(
    file_path: str, prs: PresentationType, slide_idx: int, boxes: bool = False
) -> Tuple[Path, List[str]]:
    """Rasterize one slide and digest its defects; returns image and digest.

    `boxes=False` (the QA default) renders the slide as it really looks, so the
    render answers legibility and layout questions; `boxes=True` adds the
    shape-id overlay for mapping a finding onto a box."""
    out_dir, rendered = render.render_via_soffice(
        file_path,
        out_root=Path("/tmp/pptx_render"),
        item_name="slide",
        item_idx=slide_idx,
    )
    # The soffice PDF (cached beside the rasters) carries the renderer's exact
    # word positions; reading them lets the collision check confirm overprints
    # against where glyphs actually landed rather than estimated box geometry.
    pdf_path = out_dir / f"{Path(file_path).stem}.pdf"
    slide = prs.slides[slide_idx - 1]
    raw = rendered[0]
    words = pdf_text.page_word_boxes(pdf_path, slide_idx)
    contrasts = _slide_contrasts(
        file_path, prs, slide, slide_idx, raw, words
    )
    shape_blockers = _slide_shape_blockers(file_path, prs, slide, slide_idx)
    for sid, text in _split_sentence_markers(slide):
        shape_blockers.append((
            sid,
            f"starts mid-sentence with {ellipsize(text, 40)!r}. One sentence is "
            "split across two boxes, so each reads as a fragment. Give every box "
            "a whole thought: a heading is a phrase of its own, a body is a "
            "sentence of its own",
        ))

    for sid, gap, fill in _void_markers(
        slide, _resolve_layout_chain(file_path, slide), prs.slide_height or 0
    ):
        shape_blockers.append((
            sid,
            f"leaves a {gap:.0%} band of empty slide under the title and then "
            f"fills {fill:.0%} of its box. The exemplar's boxes are placed for "
            "the copy it shipped with: clone one sized for the copy you have "
            "(the same layout with column headings above the text usually is), "
            "or move these boxes up under the title",
        ))
    for sid, pic_id, rows in pptx_contrast.text_over_artwork(
        slide, lambda shape: getattr(embedded_image(shape), "blob", None)
    ):
        shape_blockers.append((
            sid,
            f"printed over picture #{pic_id}, which carries {rows} rows of its "
            f"own text there. Clone a slide the template built to hold a title "
            "and put your copy on that, or delete this picture. Do NOT paint "
            "over the artwork or swap in a doctored copy of it - a patch shows "
            "as a flat block - and do not solve it by writing no title, which "
            "leaves the template's headline as yours",
        ))
    slide_blockers: List[str] = []
    band = rendered_void(
        slide, words, prs.slide_width or 0, prs.slide_height or 0
    )
    if band is not None:
        slide_blockers.append(
            f"a {band:.0%} band of this slide holds nothing at all. Every box "
            "on it separately fits, because each one holds less copy than the "
            "exemplar sized it for, so nothing else flags it - and the slide "
            "ships with a hole. Nudging the boxes does not close it: move them "
            "up under the title AND grow them to span the space, or clone a "
            "layout built for the amount of copy you actually have"
        )
    res = _annotate_boxes(
        raw, slide,
        prs.slide_width or 0, prs.slide_height or 0,
        effective_boxes=_text_extent_boxes(file_path, slide),
        draw_boxes=boxes,
    )
    if not res:
        return raw, _slide_findings_lines(
            slide, slide_idx, {"overlaps": [], "markers": []}, words,
            contrasts=contrasts, shape_blockers=shape_blockers,
            slide_blockers=slide_blockers,
        )
    annotated, findings = res
    return annotated, _slide_findings_lines(
        slide, slide_idx, findings, words,
        contrasts=contrasts, shape_blockers=shape_blockers,
        slide_blockers=slide_blockers,
    )


def _boxed_render(
    file_path: str,
    prs: PresentationType,
    slide_idx: int,
    render_dir: str = "/files/conversation",
    boxes: bool = False,
) -> str:
    """Render one slide + its defect digest - the diagnostic half of QA.
    Reached via --qa, not exposed on its own."""
    annotated, digest = _annotate_slide(file_path, prs, slide_idx, boxes=boxes)
    basename = os.path.splitext(os.path.basename(file_path))[0]
    published = render_publish.publish_renders(
        basename, [annotated], render_dir, _VIEW_SUBDIR
    )
    overlay = " with the #id box overlay" if boxes else ""
    lines = [f"[Slide {slide_idx} render{overlay}:]"]
    lines.extend(render_publish.render_view_lines(published))

    viewable = [scoped for _, scoped in published if scoped]
    if len(viewable) < len(published):
        # A render that is not on the conversation mount cannot be opened with
        # files__cat, so the visual check - the actual QA gate - is impossible.
        # Say so as a blocker instead of letting the pre-checks below read as a
        # pass.
        lines.append(
            "[!] some renders are NOT viewable (not on the conversation mount): "
            "the visual check cannot be performed, so QA is INCOMPLETE for "
            "those slides. Run against a deck under /files/conversation."
        )

    if digest:
        # Mechanical checks read off box geometry, the rendered pixels and the
        # PDF's word positions. They catch what a reader misses at a glance
        # (2 pt of overprint, 1.4:1 contrast) but not what only a reader sees
        # (wrong exemplar, copy that says nothing, a hole in the layout), so a
        # clean digest is a filter passed, not a slide approved.
        lines.append(
            "[Mechanical checks - a filter, not the pass. "
            "[!] fix before delivery · [w] review · [i] FYI:]"
        )
        lines.extend(digest)

    if viewable:
        lines.append(
            "[Open the render with files__cat and read it: every line legible "
            "on its background, nothing clipped, nothing bunched. "
            "--qa N --boxes relabels it with #id boxes if you need to place a "
            "finding.]"
        )
    return "\n".join(lines)


def _publish_grids(
    file_path: str,
    cells: List[Tuple[Path, str]],
    slide_nos: List[int],
    render_dir: str,
    grid_cols: int,
) -> List[Tuple[Path, Optional[str], List[int]]]:
    """Publish the grids; returns (path, scoped path, slides) each."""
    basename = os.path.splitext(os.path.basename(file_path))[0]
    # Drop missing renders here, not in the composer, so slide numbers stay aligned.
    present = [(c, n) for c, n in zip(cells, slide_nos) if c[0].exists()]
    grids = pptx_grid.compose_grids(
        [c for c, _ in present],
        Path("/tmp/pptx_render") / basename / "grids",
        grid_cols,
    )
    published = render_publish.publish_renders(
        basename, [path for path, _ in grids], render_dir, _VIEW_SUBDIR
    )
    grouped: List[Tuple[Path, Optional[str], List[int]]] = []
    consumed = 0
    for (dest, scoped), (_, labels) in zip(published, grids):
        held = [n for _, n in present[consumed : consumed + len(labels)]]
        grouped.append((dest, scoped, held))
        consumed += len(labels)
    return grouped


def print_render(
    file_path: str,
    prs: PresentationType,
    slide_nos: Optional[List[int]],
    render_dir: str = "/files/conversation",
    grid_cols: Optional[int] = None,
) -> str:
    """Plain rasterized slide(s) - no overlay - for a quick visual look. Pass a
    list of 1-based slide numbers (from a --slide N[,N,...] pattern) to rasterize
    just those in one shared pass - the deck is converted to PDF once and the PDF
    is reused across the slides - or None for the whole deck. The boxed diagnostic
    render is part of --qa, not here."""
    total_slides = len(prs.slides)
    rendered: List[Path] = []
    if slide_nos is None:
        _, rendered = render.render_via_soffice(
            file_path,
            out_root=Path("/tmp/pptx_render"),
            item_name="slide",
            item_idx=None,
        )
    else:
        for idx in slide_nos:
            if idx < 1 or idx > total_slides:
                raise ValueError(
                    f"slide index out of range: {idx} "
                    f"(deck has {total_slides} slides)"
                )
            _, one = render.render_via_soffice(
                file_path,
                out_root=Path("/tmp/pptx_render"),
                item_name="slide",
                item_idx=idx,
            )
            rendered.extend(one)

    rendered_nos = slide_nos or list(range(1, len(rendered) + 1))
    if grid_cols:
        grids = _publish_grids(
            file_path,
            [(p, f"slide {n}") for p, n in zip(rendered, rendered_nos)],
            rendered_nos,
            render_dir,
            grid_cols,
        )
        lines = [
            f"[Rendered {len(rendered)} slides into {len(grids)} grid(s), "
            f"{grid_cols}/row:]"
        ]
        lines.extend(pptx_grid.grid_lines(grids))
        return "\n".join(lines)

    basename = os.path.splitext(os.path.basename(file_path))[0]
    published = render_publish.publish_renders(
        basename, rendered, render_dir, _VIEW_SUBDIR
    )
    plural = "" if len(published) == 1 else "s"
    lines = [f"[Rendered {len(published)} slide{plural}:]"]
    lines.extend(render_publish.render_view_lines(published))
    return "\n".join(lines)


def print_qa(
    file_path: str,
    prs: PresentationType,
    slide_nos: List[int],
    render_dir: str = "/files/conversation",
    boxes: bool = False,
) -> str:
    """QA gate - run after a round of edits. Accepts one slide or several (a
    pattern like `2,5,7-9`); for each, bundles the slide's authoritative text
    (#id-tagged, to read back) with a full-size render of the slide and its
    defect digest. QA-ing several slides at once shares a single soffice
    conversion (the PDF is cached after the first render), so batching the
    changed slides is much faster than one call per slide.

    One image per slide, deliberately. Tiling the batch into a grid quarters
    each slide's pixels, and a slide read at a quarter size hides exactly the
    defects QA exists to catch: a caption the same colour as its background, a
    word broken across a pill, half an inch of a box past the slide edge."""
    total_slides = len(prs.slides)
    for slide_idx in slide_nos:
        if slide_idx < 1 or slide_idx > total_slides:
            raise ValueError(
                f"slide index out of range: {slide_idx} "
                f"(deck has {total_slides} slides)"
            )
    blocks = [
        f"[QA slide {slide_idx} - #id text + render path:]\n\n"
        + print_text(prs, slide_idx)
        + "\n\n"
        + _boxed_render(file_path, prs, slide_idx, render_dir, boxes=boxes)
        for slide_idx in slide_nos
    ]
    return "\n\n".join(blocks)


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="pptx_inspect",
        usage=USAGE,
        add_help=False,
    )
    parser.add_argument("file", nargs="?")
    parser.add_argument("--slide", metavar="N[,N,...]")
    parser.add_argument("--layouts", action="store_true")
    parser.add_argument("--text", action="store_true")
    parser.add_argument("--media", action="store_true")
    parser.add_argument("--render", action="store_true")
    parser.add_argument("--qa", metavar="N[,N,...]")
    parser.add_argument("--grid", action="store_true")
    parser.add_argument("--boxes", action="store_true")
    parser.add_argument(
        "--grid-cols",
        dest="grid_cols",
        type=int,
        default=pptx_grid.DEFAULT_GRID_COLS,
        metavar="N",
    )
    parser.add_argument(
        "--render-dir",
        dest="render_dir",
        metavar="DIR",
        default="/files/conversation",
    )
    parser.add_argument("--compare")
    parser.add_argument("--validate", action="store_true")
    parser.add_argument("--max-shapes", type=int, default=DEFAULT_MAX_SHAPES)
    parser.add_argument("--offset", type=int, default=0)
    parser.add_argument("--help", "-h", action="store_true", dest="help_flag")

    args = parser.parse_args()

    if args.help_flag:
        sys.stdout.write(HELP_TEXT + "\n")
        return 0

    if not args.file:
        sys.stderr.write(f"Error: file is required\nUsage: {USAGE}\n")
        return 1

    if not os.path.isfile(args.file):
        sys.stderr.write(f"Error: file not found: {args.file}\n")
        return 1

    if args.max_shapes < 1:
        sys.stderr.write("Error: --max-shapes must be >= 1\n")
        return 1
    if args.offset < 0:
        sys.stderr.write("Error: --offset must be >= 0\n")
        return 1
    if not 1 <= args.grid_cols <= pptx_grid.MAX_GRID_COLS:
        sys.stderr.write(
            f"Error: --grid-cols must be 1..{pptx_grid.MAX_GRID_COLS}\n"
        )
        return 1
    grid_cols = args.grid_cols if args.grid else None
    if args.compare is not None and not os.path.isfile(args.compare):
        sys.stderr.write(f"Error: --compare file not found: {args.compare}\n")
        return 1

    file_header = (
        f"[File: {os.path.basename(args.file)} | "
        f"{format_size(os.path.getsize(args.file))}]"
    )

    if args.validate:
        body = pptx_validate.report(args.file, args.compare)
    elif args.media:
        body = print_media(args.file)
    elif args.compare is not None:
        body = print_compare(args.file, args.compare)
    else:
        prs = Presentation(args.file)
        if args.qa is not None:
            body = print_qa(
                args.file,
                prs,
                parse_slide_patterns(args.qa),
                args.render_dir,
                boxes=args.boxes,
            )
            if grid_cols:
                body = (
                    "[--grid is for --render browsing; QA renders one slide per "
                    "image so nothing is too small to read.]\n" + body
                )
        elif args.render:
            slide_nos = parse_slide_patterns(args.slide) if args.slide else None
            body = print_render(
                args.file, prs, slide_nos, args.render_dir, grid_cols
            )
        elif args.layouts:
            body = print_layouts(prs, args.file)
        elif args.text:
            body = print_text(prs)
        elif args.slide is not None:
            body = print_slides(
                prs, args.file, parse_slide_patterns(args.slide),
                args.offset, args.max_shapes,
            )
        else:
            body = print_overview(prs, args.file)

    full = file_header + "\n" + body
    text, truncated = safe_output(full)
    sys.stdout.write(text)
    sys.stdout.write("\n")
    if truncated:
        sys.stdout.write(
            "[Output truncated; narrow with --slide or paginate with "
            "--offset / --max-shapes]\n"
        )
    return 0


if __name__ == "__main__":
    try:
        sys.exit(main())
    except (ValueError, KeyError) as exc:
        sys.stderr.write(f"Error: {exc}\n")
        sys.exit(1)
    except Exception as exc:
        sys.stderr.write(f"Error: {type(exc).__name__}: {exc}\n")
        sys.exit(1)
