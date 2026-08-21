"""Geometry primitives for pptx_inspect: EMU/box math, the text-fit estimate,
and the overlap classifier, plus the shared geometry thresholds.

The no-dependency layer of the inspector - pure box math over python-pptx shapes,
importing nothing from the other pptx_* modules - so audit/render can build on it
without an import cycle.
"""
from __future__ import annotations

from typing import NamedTuple, Optional, Tuple

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.base import BaseShape

import pptx_fontmetrics

EMU_PER_INCH = 914_400
EDGE_EPSILON_EMU = 45_720  # 0.05" tolerance before flagging edge overflow.

# Text-fit estimation. What decides whether text fits is the box geometry at
# the chosen font size, not the word count - so we estimate how many characters
# a box holds and surface it, to stop the agent sizing text blindly. It is a
# rough ESTIMATE (mean advance width over the string, no real line breaking),
# so the overset warning only fires on GROSS overflow of a genuine
# multi-line container, never on a box short enough to be a nominal-height label
# (text overflows those by design) or one set to grow to fit its text.
CHAR_WIDTH_EM = 0.5  # avg proportional Latin glyph advance, in em
LINE_HEIGHT_FACTOR = 1.2  # typical single line height, in em
TEXTBOX_MARGIN_W_IN = 0.2  # PowerPoint default internal left+right inset
TEXTBOX_MARGIN_H_IN = 0.1  # PowerPoint default internal top+bottom inset
OVERSET_TOLERANCE = 1.8  # flag only when text exceeds capacity by this factor
FILL_FLOATING = 0.30  # text using less than this fraction of a real box "floats"
# Rendered-extent estimate (grows a box to wrap copy that overflows it), biased
# larger so the grown box overshoots the text rather than under-covering it and
# missing an overlap. Substitute fonts (the brand face is absent at render time)
# run wider AND taller than the fit estimate's 0.5em / 1.2, so the extent uses
# its own, larger constants for BOTH how many lines the copy wraps to and how
# tall each line is. Height availability still uses LINE_HEIGHT_FACTOR above, so
# a real multi-line overflow grows the box, and the growth is capped.
EXTENT_CHAR_WIDTH_EM = 0.6  # wider glyph advance for the extent wrap estimate
EXTENT_LINE_HEIGHT_FACTOR = 1.4  # generous line height for the extent height
EXTENT_PAD_IN = 0.1  # extra height added to the extent estimate
EXTENT_OVERFLOW_SLACK = 1  # grow when copy needs >= this many more lines than fit
EXTENT_MAX_GROWTH = 2.5  # cap: grown height <= this multiple of the declared

ASPECT_TOLERANCE = 1.18  # flag image stretch/squish beyond ~15%
MIN_IMAGE_DPI = 70  # flag pictures displayed below this effective resolution
# Overlap classification. Overlap is often intentional (a label on a photo, a
# caption on a pill), so raw intersection is not a defect. We separate three
# cases by per-box containment (intersection / each box's area):
#   - both boxes >= STACKED_CONT covered  -> the boxes coincide (STACKED): a
#     real bug (e.g. a title and body dropped at the same spot) -> [!] blocker.
#   - one box >= CONTAINMENT_TAU inside the other -> foreground on background
#     (caption on a photo, text on a pill): intentional -> suppressed.
#   - otherwise, shallower-axis overlap >= PEER_PENETRATION_EMU -> a partial
#     peer overlap: could be a designed overlay or an accidental collision, so
#     it is surfaced as a quantitative [i] ADVISORY (judge it in the render),
#     never a hard blocker - that would fail the template's own overlays.
STACKED_CONT = 0.90  # both boxes >= this covered => effectively the same box
CONTAINMENT_TAU = 0.85  # one box >= this fraction inside the other => fg/bg, ok
PEER_PENETRATION_EMU = EMU_PER_INCH // 10  # 0.1" shallow-axis overlap to note it
CENTER_OFFSET_EMU = int(EMU_PER_INCH * 0.2)  # inner box off-center beyond this => advisory
SAFE_MARGIN_EMU = int(EMU_PER_INCH * 0.2)  # text crowding within this of an edge => advisory
FULL_SPAN = 0.9  # a shape covering >= this fraction of a slide axis is a banner/background


def emu_to_inches(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    return emu / EMU_PER_INCH


def format_box(shape: BaseShape) -> str:
    left = emu_to_inches(shape.left)
    top = emu_to_inches(shape.top)
    width = emu_to_inches(shape.width)
    height = emu_to_inches(shape.height)
    if None in (left, top, width, height):
        return "(?,?)"
    return f"({left:.1f},{top:.1f}) {width:.1f}x{height:.1f}\""


class FitEstimate(NamedTuple):
    chars_per_line: int
    capacity: Optional[int]  # total chars; None when height isn't a real constraint
    measured: bool = False  # widths read off the real face, not the constant


def _frame_text_len(shape: BaseShape) -> int:
    if not shape.has_text_frame:
        return 0
    return sum(len(p.text or "") for p in shape.text_frame.paragraphs)


def _fit_estimate(
    shape: BaseShape, size_pt: float, typeface: Optional[str] = None
) -> Optional[FitEstimate]:
    if shape.width is None or shape.height is None or size_pt <= 0:
        return None
    w_in = shape.width / EMU_PER_INCH
    h_in = shape.height / EMU_PER_INCH
    text = shape.text_frame.text if getattr(shape, "has_text_frame", False) else ""
    char_em, measured = pptx_fontmetrics.mean_char_em(text, typeface)
    char_w = size_pt * char_em / 72.0
    line_h = size_pt * LINE_HEIGHT_FACTOR / 72.0
    if char_w <= 0 or line_h <= 0:
        return None
    cpl = int(max(0.0, w_in - TEXTBOX_MARGIN_W_IN) / char_w)
    lines = int(max(0.0, h_in - TEXTBOX_MARGIN_H_IN) / line_h)
    if cpl < 1:
        return None
    # Height constrains only genuine multi-line containers; a box shorter than
    # ~2 lines is a nominal-height label whose text overflows by design.
    capacity = cpl * lines if lines >= 2 else None
    return FitEstimate(cpl, capacity, measured)


def _grows_to_fit(shape: BaseShape) -> bool:
    return (
        shape.has_text_frame
        and shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    )


def _text_extent_box(
    shape: BaseShape, size_pt: float
) -> Optional[Tuple[int, int, int, int]]:
    """The box grown to wrap this shape's rendered copy when the text needs more
    lines than the declared box holds (a multi-paragraph or wrapped block in a
    too-short box, like a two-line title in a one-line box). Returns
    (left, top, width, height) in EMU - kept at the declared top and grown
    DOWNWARD (a BOTTOM-anchored box, whose text spills up, grown upward) - or
    None when the text fits, so callers fall back to the declared box.

    Growth triggers on a genuine multi-line overflow (>= EXTENT_OVERFLOW_SLACK
    more lines than the box holds), is biased larger so the box overshoots
    rather than under-covers the text (EXTENT_*), and is capped at
    EXTENT_MAX_GROWTH x the declared height so a mis-estimated wrap can't balloon
    it across the slide. The cap, not the trigger, is what bounds over-growth, so
    a real one-line title overflow is still surfaced. A MIDDLE-anchored box is
    never grown symmetrically."""
    if not shape.has_text_frame:
        return None
    if None in (shape.left, shape.top, shape.width, shape.height):
        return None
    if size_pt <= 0:
        return None
    w_in = shape.width / EMU_PER_INCH
    h_in = shape.height / EMU_PER_INCH
    paragraphs = list(shape.text_frame.paragraphs)
    last = -1
    for i, p in enumerate(paragraphs):
        if (p.text or "").strip():
            last = i
    if last < 0:
        return None

    # Lines the copy needs at a given glyph width, summed over paragraphs (each
    # wraps to ceil(chars/cpl); an empty interior spacer still occupies one row;
    # trailing empties ignored). None when the box is too narrow to hold a char.
    def _lines_for(char_em: float) -> Optional[int]:
        cw = size_pt * char_em / 72.0
        if cw <= 0:
            return None
        cpl = int(max(0.0, w_in - TEXTBOX_MARGIN_W_IN) / cw)
        if cpl < 1:
            return None
        return sum(
            max(1, -(-len(p.text or "") // cpl))  # ceil division
            for p in paragraphs[: last + 1]
        )

    # Two widths, two jobs. The conservative width (0.5em, the same as
    # _fit_estimate) decides WHETHER the copy overflows: borderline copy that
    # only the wider extent width would wrap is left alone, so a one-line title
    # that just fits never fabricates a spill onto whatever sits below it. The
    # generous width (0.6em - substitute fonts run wide) then decides HOW FAR to
    # grow, so a genuine overflow still overshoots the text rather than
    # under-covering it and missing a real collision.
    lines_trigger = _lines_for(CHAR_WIDTH_EM)
    lines_extent = _lines_for(EXTENT_CHAR_WIDTH_EM)
    if lines_trigger is None or lines_extent is None:
        return None
    line_h = size_pt * LINE_HEIGHT_FACTOR / 72.0
    if line_h <= 0:
        return None
    lines_avail = int(max(0.0, h_in - TEXTBOX_MARGIN_H_IN) / line_h)
    # A genuine multi-line overflow grows the box: the copy must need at least
    # EXTENT_OVERFLOW_SLACK more lines than the box holds and span at least two
    # lines. A single line fits (or is a nominal-height label that overflows by
    # design) and is never grown. One extra line is real overflow - a two-line
    # title in a one-line box spills onto whatever sits below it - so it grows;
    # EXTENT_MAX_GROWTH below bounds how far a mis-estimated wrap can take it.
    if lines_trigger < 2 or lines_trigger < lines_avail + EXTENT_OVERFLOW_SLACK:
        return None
    gen_line_h = size_pt * EXTENT_LINE_HEIGHT_FACTOR / 72.0
    ext_h_in = lines_extent * gen_line_h + TEXTBOX_MARGIN_H_IN + EXTENT_PAD_IN
    new_h = int(round(ext_h_in * EMU_PER_INCH))
    # Cap the growth so a mis-estimated wrap can't balloon the box across the
    # slide (it still overshoots the text, just not unboundedly).
    new_h = min(new_h, int(round(shape.height * EXTENT_MAX_GROWTH)))
    if new_h <= shape.height:
        return None
    # Keep the declared top and grow DOWNWARD for TOP/MIDDLE/inherited anchors: a
    # MIDDLE box is never grown symmetrically, which used to drift it up off its
    # own position. Only a BOTTOM-anchored box, whose text spills upward, grows
    # up - clamped to the slide top.
    va = getattr(shape.text_frame, "vertical_anchor", None)
    va_name = (getattr(va, "name", None) or "TOP").upper()
    left, top, height = shape.left, shape.top, shape.height
    if va_name == "BOTTOM":
        new_top = max(0, top + height - new_h)
    else:
        new_top = top
    return (left, new_top, shape.width, new_h)


def _classify_overlap(a: Tuple[int, int, int, int],
                      b: Tuple[int, int, int, int]):
    """Classify the overlap of two (left, top, width, height) boxes.
    Returns (kind, penetration_emu, axis):
      - ("stacked", 0, "")     both boxes >= STACKED_CONT covered (coincide)
      - ("contained", 0, "")   one box >= CONTAINMENT_TAU inside the other (fg/bg)
      - ("peer", pen, axis)    partial overlap >= PEER_PENETRATION_EMU on the
                               shallow axis (a designed overlay or a collision)
      - (None, 0, "")          no/negligible overlap
    Shared by the --slide text lint and the --render overlap shading so both
    agree on what counts as a collision."""
    al, at, aw, ah = a
    bl, bt, bw, bh = b
    if min(aw, ah, bw, bh) <= 0:
        return (None, 0, "")
    ix = min(al + aw, bl + bw) - max(al, bl)
    iy = min(at + ah, bt + bh) - max(at, bt)
    if ix <= 0 or iy <= 0:
        return (None, 0, "")
    inter = ix * iy
    cont_a = inter / (aw * ah)
    cont_b = inter / (bw * bh)
    if min(cont_a, cont_b) >= STACKED_CONT:
        return ("stacked", 0, "")
    if max(cont_a, cont_b) >= CONTAINMENT_TAU:
        return ("contained", 0, "")
    if min(ix, iy) >= PEER_PENETRATION_EMU:
        return ("peer", min(ix, iy), "horizontally" if ix < iy else "vertically")
    return (None, 0, "")


def _render_collision(
    decl_a: Tuple[int, int, int, int],
    decl_b: Tuple[int, int, int, int],
    eff_a: Tuple[int, int, int, int],
    eff_b: Tuple[int, int, int, int],
) -> Tuple[bool, int, str]:
    """Whether the boxed render should flag two shapes as colliding, judged on
    their effective (text-extent) boxes, plus the shallow-axis penetration (EMU)
    and a kind the QA digest tiers on.

    A peer overlap counts, as before. So does a containment that exists ONLY
    after a box grew to wrap overflowing text - that text spilled onto a
    neighbour. But a containment that already holds at the DECLARED sizes is an
    intentional foreground/background overlay (a label on a card) and does not.
    Returns (flag, penetration_emu, kind) where kind is "spill" (a box overflowed
    its declared bounds onto a neighbour - almost always a defect), "peer" (a
    partial overlap already present at declared geometry - often a designed
    overlay), or "" when there is no collision."""
    kind, pen, _ = _classify_overlap(eff_a, eff_b)
    if kind == "peer":
        return (True, pen, "peer")
    if kind == "contained":
        if _classify_overlap(decl_a, decl_b)[0] == "contained":
            return (False, 0, "")  # designed fg/bg overlay, not a collision
        ix = min(eff_a[0] + eff_a[2], eff_b[0] + eff_b[2]) - max(eff_a[0], eff_b[0])
        iy = min(eff_a[1] + eff_a[3], eff_b[1] + eff_b[3]) - max(eff_a[1], eff_b[1])
        return (True, min(ix, iy), "spill")  # spillover depth on the shallow axis
    return (False, 0, "")


def shape_kind(shape: BaseShape) -> str:
    if shape.has_chart:
        return "chart"
    if shape.has_table:
        return "table"
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "pic"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "ph"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto"
    if st == MSO_SHAPE_TYPE.LINE:
        return "line"
    if st == MSO_SHAPE_TYPE.FREEFORM:
        return "free"
    if st == MSO_SHAPE_TYPE.MEDIA:
        return "media"
    name = getattr(st, "name", None)
    return name.lower() if name else "shape"
