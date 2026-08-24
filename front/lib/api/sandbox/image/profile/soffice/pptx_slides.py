#!/opt/venv/bin/python3
"""pptx_slides - safe slide-level structural edits for .pptx decks.

python-pptx has no slide duplicate / move / reorder API, and the common
community recipes corrupt the package (two slides sharing one chart part,
stale relationship ids). This tool does those operations correctly so an
agent editing a template in place can add, move, and remove slides without
breaking it.

Operations (one per call, 1-indexed slide numbers, edits the file in place):
  --duplicate P [--count K] [--after M]   copy the slides in pattern P (each K
                                          times), placed after slide M (default:
                                          after the last slide in P)
  --move N --to M                         move slide N to position M
  --delete P                              delete the slides in pattern P,
                                          dropping only their own relationships

A slide pattern P is a comma-separated list of slide numbers and inclusive
ranges, e.g. `5`, `2,5,8`, `3-7`, or `2,5,7-9`; duplicates clone in the order
given.

Duplication shares immutable image/media parts (correct - that is what
PowerPoint does) and deep-clones mutable parts (charts and their embedded
worksheets, SmartArt/diagrams, OLE objects) so each copy is independent.
Relationship ids in the copied XML are rewritten to the new ids. Speaker
notes are not carried onto a duplicate.
"""

from __future__ import annotations

import argparse
import copy
import os
import re
import sys
from typing import Dict, List, Optional

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.oxml.ns import qn

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Relationship types whose target part is immutable and safe to share between
# slides (the same bytes referenced from several places). Everything else that
# is internal gets deep-cloned so each duplicate owns its copy.
SHAREABLE = frozenset(
    {RT.IMAGE, RT.MEDIA, getattr(RT, "VIDEO", "__novideo__"),
     getattr(RT, "AUDIO", "__noaudio__")}
)
# Slide relationships that must not be carried onto a duplicate: the layout
# (the new slide gets its own) and the notes slide (back-references the source).
SKIP_ON_SLIDE = frozenset({RT.SLIDE_LAYOUT, RT.NOTES_SLIDE})

USAGE = (
    "pptx_slides <file> (--duplicate N[,N,...] [--count K] [--after M] "
    "| --move N --to M | --delete N[,N,...])"
)


def _partname_template(partname: str) -> str:
    """'/ppt/charts/chart1.xml' -> '/ppt/charts/chart%d.xml' so next_partname
    can allocate a fresh, unused name of the same kind."""
    s = str(partname)
    new, n = re.subn(r"(\d+)(\.[^.]+)$", r"%d\2", s)
    if n:
        return new
    # No trailing number before the extension: insert one.
    return re.sub(r"(\.[^.]+)$", r"%d\1", s)


def _rewrite_rids(root: etree._Element, rid_map: Dict[str, str]) -> None:
    """Rewrite every relationship-namespace attribute (r:id, r:embed, r:link,
    r:dm, …) whose value is a remapped id, in place across the element tree."""
    if not rid_map:
        return
    for el in root.iter():
        for name, val in list(el.attrib.items()):
            if name.startswith("{" + R_NS + "}") and val in rid_map:
                el.set(name, rid_map[val])


def _clone_part(part, pkg) -> Part:
    """Deep-clone a part and the sub-tree of mutable parts it depends on
    (sharing immutable image/media parts), rewriting relationship ids inside
    XML parts so the clone is internally consistent."""
    new_partname = pkg.next_partname(_partname_template(part.partname))
    new_part = Part(new_partname, part.content_type, pkg, part.blob)

    rid_map: Dict[str, str] = {}
    for rId, rel in part.rels.items():
        if rel.is_external:
            new_rid = new_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            target = rel.target_part
            child = target if rel.reltype in SHAREABLE else _clone_part(target, pkg)
            new_rid = new_part.relate_to(child, rel.reltype)
        rid_map[rId] = new_rid

    if rid_map and str(part.content_type).endswith("+xml"):
        root = etree.fromstring(part.blob)
        _rewrite_rids(root, rid_map)
        new_part._blob = etree.tostring(
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        )
    return new_part


def _copy_slide_rels(src_slide, dest_slide, pkg) -> Dict[str, str]:
    """Recreate the source slide's relationships on the destination slide
    (sharing images, deep-cloning charts/diagrams/OLE, copying externals),
    returning a map of old relationship id -> new relationship id."""
    rid_map: Dict[str, str] = {}
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype in SKIP_ON_SLIDE:
            continue
        if rel.is_external:
            new_rid = dest_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
        else:
            target = rel.target_part
            child = target if rel.reltype in SHAREABLE else _clone_part(target, pkg)
            new_rid = dest_slide.part.relate_to(child, rel.reltype)
        rid_map[rId] = new_rid
    return rid_map


def _clone_slide_append(prs, src_index: int):
    """Duplicate the slide at 0-based src_index, appending the copy at the end of
    the slide order, and return its new `sldId` element. Appending (rather than
    inserting at a target) keeps every existing slide's index stable while a
    batch of copies is produced; the caller places the copies afterwards."""
    src = prs.slides[src_index]
    pkg = prs.part.package

    dest = prs.slides.add_slide(src.slide_layout)  # appended; carries layout rel
    # Drop the auto-cloned layout placeholders; we copy the source's own shapes.
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    # Copy a slide background override, if the source sets one.
    src_bg = src._element.cSld.find(qn("p:bg"))
    if src_bg is not None and dest._element.cSld.find(qn("p:bg")) is None:
        dest._element.cSld.insert(0, copy.deepcopy(src_bg))

    sp_tree = dest.shapes._spTree
    for shape in src.shapes:
        sp_tree.append(copy.deepcopy(shape._element))

    rid_map = _copy_slide_rels(src, dest, pkg)
    # Rewrite across the whole cSld, not just the shape tree: a slide background
    # override lives in <p:bg>, a sibling of <p:spTree>, and its picture fill
    # carries an r:embed of its own. Rewriting only the shapes left that r:embed
    # pointing at the SOURCE slide's relationship id, which on the copy resolves
    # to whatever else happens to hold that id - so duplicating a slide with a
    # picture background silently gave the copy a different image behind it.
    _rewrite_rids(dest._element.cSld, rid_map)

    return list(prs.slides._sldIdLst)[-1]  # add_slide appended it last


def op_duplicate(
    prs, slide_nos: List[int], count: int, after: Optional[int]
) -> str:
    n = len(prs.slides)
    for s in slide_nos:
        if not 1 <= s <= n:
            raise ValueError(f"--duplicate {s} out of range (1..{n})")
    if count < 1:
        raise ValueError("--count must be >= 1")
    # Default placement keeps the copies next to the originals they came from.
    after_no = max(slide_nos) if after is None else after
    if not 0 <= after_no <= n:
        raise ValueError(f"--after {after_no} out of range (0..{n})")

    sld_id_lst = prs.slides._sldIdLst
    # Resolve the anchor to a stable sldId element before any insertion shifts
    # indices; None means "front of the deck" (--after 0).
    anchor = list(sld_id_lst)[after_no - 1] if after_no >= 1 else None

    # Clone every requested slide (each `count` times) to the end first - which
    # leaves all source indices stable - then move the whole batch, in order, to
    # its destination.
    new_ids = [
        _clone_slide_append(prs, s - 1)
        for s in slide_nos
        for _ in range(count)
    ]
    for nid in new_ids:
        sld_id_lst.remove(nid)
    pos = list(sld_id_lst).index(anchor) + 1 if anchor is not None else 0
    for offset, nid in enumerate(new_ids):
        sld_id_lst.insert(pos + offset, nid)

    where = (
        "end" if after_no == n
        else "front" if after_no == 0
        else f"after slide {after_no}"
    )
    pretty = ",".join(map(str, slide_nos))
    return (
        f"Duplicated slide(s) {pretty} x{count} ({where}); "
        f"deck now has {len(prs.slides)} slides"
    )


def op_move(prs, slide_no: int, to: int) -> str:
    n = len(prs.slides)
    if not 1 <= slide_no <= n:
        raise ValueError(f"--move {slide_no} out of range (1..{n})")
    if not 1 <= to <= n:
        raise ValueError(f"--to {to} out of range (1..{n})")
    sld_id_lst = prs.slides._sldIdLst
    sld_id = list(sld_id_lst)[slide_no - 1]
    sld_id_lst.remove(sld_id)
    sld_id_lst.insert(to - 1, sld_id)
    return f"Moved slide {slide_no} to position {to}"


def _purge_slide_refs(pres_el: etree._Element, rId: str) -> None:
    """Drop presentation.xml elements still pointing at `rId`.

    A custom show reuses the slide's rId, and drop_rel keeps any rel still
    referenced - so the slide leaves the order but its part stays as an orphan.
    """
    for el in list(pres_el.iter()):
        if any(
            name.startswith("{" + R_NS + "}") and value == rId
            for name, value in el.attrib.items()
        ):
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)


def op_delete(prs, slide_nos: List[int]) -> str:
    n = len(prs.slides)
    for s in slide_nos:
        if not 1 <= s <= n:
            raise ValueError(f"--delete {s} out of range (1..{n})")
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    # Delete from the highest index down so earlier indices stay valid.
    for s in sorted(set(slide_nos), reverse=True):
        sld_id = sld_ids[s - 1]
        rId = sld_id.get(qn("r:id"))
        sld_id_lst.remove(sld_id)
        _purge_slide_refs(prs.part._element, rId)
        prs.part.drop_rel(rId)  # drop only this slide's rel; shared media stays
        if rId in prs.part.rels:
            raise ValueError(
                f"slide {s} is still referenced from presentation.xml after "
                f"its slide order entry was removed ({rId}); deleting it would "
                "strand its part in the package"
            )
    return f"Deleted slide(s) {','.join(map(str, sorted(set(slide_nos))))}"


def _parse_slide_patterns(raw: str) -> List[int]:
    """Expand a slide pattern (comma-separated numbers and inclusive `A-B`
    ranges, e.g. `2,5,7-9`) into an ordered, de-duplicated list of 1-based slide
    numbers. Order is preserved so `--duplicate` clones in the order given."""
    out: List[int] = []
    seen = set()

    def add(value: int) -> None:
        if value not in seen:
            seen.add(value)
            out.append(value)

    for token in raw.split(","):
        token = token.strip()
        if not token:
            continue
        if "-" in token:
            lo_s, _, hi_s = token.partition("-")
            lo_s, hi_s = lo_s.strip(), hi_s.strip()
            if not (lo_s.isdigit() and hi_s.isdigit()):
                raise ValueError(f"invalid slide range: {token!r}")
            lo, hi = int(lo_s), int(hi_s)
            if lo > hi:
                raise ValueError(f"invalid slide range (start > end): {token!r}")
            for value in range(lo, hi + 1):
                add(value)
        elif token.isdigit():
            add(int(token))
        else:
            raise ValueError(f"invalid slide number: {token!r}")
    if not out:
        raise ValueError("no slides specified")
    return out


def main() -> int:
    parser = argparse.ArgumentParser(
        prog="pptx_slides",
        usage=USAGE,
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
        add_help=True,
    )
    parser.add_argument("file", nargs="?")
    group = parser.add_mutually_exclusive_group()
    group.add_argument("--duplicate", metavar="N[,N,...]")
    group.add_argument("--move", type=int, metavar="N")
    group.add_argument("--delete", metavar="N[,N,...]")
    parser.add_argument("--count", type=int, default=1)
    parser.add_argument("--after", type=int)
    parser.add_argument("--to", type=int)
    args = parser.parse_args()

    if not args.file:
        sys.stderr.write(f"Error: file is required\nUsage: {USAGE}\n")
        return 1
    if not os.path.isfile(args.file):
        sys.stderr.write(f"Error: file not found: {args.file}\n")
        return 1

    prs = Presentation(args.file)
    try:
        if args.duplicate is not None:
            msg = op_duplicate(
                prs, _parse_slide_patterns(args.duplicate), args.count, args.after
            )
        elif args.move is not None:
            if args.to is None:
                raise ValueError("--move requires --to")
            msg = op_move(prs, args.move, args.to)
        elif args.delete is not None:
            msg = op_delete(prs, _parse_slide_patterns(args.delete))
        else:
            sys.stderr.write(f"Error: choose an operation\nUsage: {USAGE}\n")
            return 1
    except ValueError as exc:
        sys.stderr.write(f"Error: {exc}\n")
        return 1

    prs.save(args.file)
    sys.stdout.write(msg + "\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
