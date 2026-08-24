"""Bounding-box render overlay for pptx_inspect: draw each shape's box on the
rasterized slide, detect rendered text-row positions by edge density, and flag
decorative-marker runs left with no text row beside them.

The top layer of the visual QA - depends on geometry (shape_kind,
_classify_overlap, EMU) and PIL; the CLI calls _annotate_boxes from --qa.
"""
from __future__ import annotations

from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from pptx_geometry import EMU_PER_INCH, FULL_SPAN, _render_collision, shape_kind

# (left, top, width, height) in EMU.
BoxEmu = Tuple[int, int, int, int]


def _is_full_span(box: BoxEmu, slide_w_emu: int, slide_h_emu: int) -> bool:
    """A shape covering >= FULL_SPAN of either slide axis is a banner or
    background: its overlap with foreground content is intentional layering, not
    a collision, so it is excluded from the overlap wash."""
    _, _, w, h = box
    return w >= FULL_SPAN * slide_w_emu or h >= FULL_SPAN * slide_h_emu


def _rects_overlap(a, b) -> bool:
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _suppress_as_layering(
    a: BoxEmu, b: BoxEmu, slide_w_emu: int, slide_h_emu: int
) -> bool:
    """Whether a full-span shape makes this pair's overlap intentional layering
    (content sitting on a banner/backdrop) rather than a collision.

    Boxes are (left, top, width, height) in EMU. A full-span shape suppresses
    the pair ONLY when the two boxes already overlap at their declared geometry.
    A box that grew past its declared bounds and newly spilled onto a full-span
    shape (e.g. an overflowing title onto a full-width subtitle) is a real
    overflow, not layering, so it is not suppressed here - the caller's
    collision test surfaces it."""
    if not (
        _is_full_span(a, slide_w_emu, slide_h_emu)
        or _is_full_span(b, slide_w_emu, slide_h_emu)
    ):
        return False
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (
        ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay
    )


def _contrast_text(color):
    """Black or white text, whichever reads on a chip of `color`."""
    r, g, b = color[:3]
    lum = 0.299 * r + 0.587 * g + 0.114 * b
    return (0, 0, 0, 255) if lum > 140 else (255, 255, 255, 255)


# Text-row detection by horizontal EDGE density. A box that spans a gradient or
# blob background has a constant colour-vs-median "ink" baseline that collapses
# to one fake row, but glyph strokes still produce many sharp horizontal
# transitions while smooth backgrounds produce ~none - so edge count separates
# text from background regardless of contrast polarity or gradient. Returns each
# detected line's vertical centre (px): the input for marker alignment.
#
# Contrast is NOT measured here. Sampling it per detected row was tried and
# dropped - a row strip mixes glyphs with whatever the box sits on, so legible
# white-on-brand text scored as low as broken dark-on-dark. `pptx_contrast`
# measures it per rendered WORD box instead, where the ink is a solid fraction of
# the pixels, and that separates the two cleanly (1.1-1.6 broken, 2.4+ legible).
TEXTROW_EDGE_DELTA = 60  # adjacent-sample manhattan diff that marks a glyph edge
TEXTROW_X_STEP = 2  # sample every Nth column (speed; detection is robust to it)


def _text_row_centers(
    rgb_image, box_px: Tuple[float, float, float, float]
) -> List[float]:
    x0, y0, x1, y1 = (int(box_px[0]), int(box_px[1]),
                      int(box_px[2]), int(box_px[3]))
    x0, y0 = max(0, x0), max(0, y0)
    x1, y1 = min(rgb_image.width, x1), min(rgb_image.height, y1)
    if x1 - x0 < 6 or y1 - y0 < 6:
        return []
    px = rgb_image.load()
    xs = list(range(x0, x1, TEXTROW_X_STEP))
    edge_min = max(8, len(xs) // 15)  # text rows show many edges; background ~0-3

    def edge_count(y):
        line = [px[x, y] for x in xs]
        return sum(
            1 for a, b in zip(line, line[1:])
            if abs(a[0] - b[0]) + abs(a[1] - b[1]) + abs(a[2] - b[2])
            > TEXTROW_EDGE_DELTA
        )

    out: List[float] = []
    group: List[int] = []

    def flush():
        if len(group) >= 2:  # a text line spans more than one pixel row
            out.append((min(group) + max(group)) / 2.0)

    for y in range(y0, y1):
        if edge_count(y) >= edge_min:
            group.append(y)
        else:
            flush()
            group = []
    flush()
    return out


_BOX_PALETTE = [
    (255, 64, 64), (64, 160, 255), (0, 200, 120), (255, 170, 0),
    (200, 80, 255), (0, 200, 200), (255, 90, 170), (150, 210, 0),
]

# Decorative row-markers (checkmarks, bullets, numbers, icons) are placed at
# FIXED positions to sit beside specific text rows. When filled copy wraps, the
# rendered rows drift off those positions and a marker strands in empty space -
# a defect box-geometry can't see (the marker didn't move; the text reflowed).
# Detect a marker as a small shape that forms a regular vertical run of >=3
# beside a text box, then check each one has a rendered text row near it.
MARKER_MAX_DIM_EMU = int(EMU_PER_INCH * 0.6)  # a row-marker is small both ways
MARKER_X_TOL_EMU = int(EMU_PER_INCH * 0.15)  # same-column left tolerance
MARKER_MIN_RUN = 3  # a run is at least this many aligned markers
MARKER_PAIR_GAP_EMU = EMU_PER_INCH  # max horizontal gap to pair a run to a text box
MARKER_ALIGN_TOL_IN = 0.15  # a marker within this of a text row counts as aligned


def _marker_runs(shapes: List[BaseShape]) -> List[List[BaseShape]]:
    """Group small decorative shapes into vertical runs (>=3 sharing a left
    edge). Each run is returned sorted top-to-bottom."""
    cands = sorted(
        [
            s for s in shapes
            if shape_kind(s) in ("pic", "auto")
            and s.width and s.height
            and s.width <= MARKER_MAX_DIM_EMU and s.height <= MARKER_MAX_DIM_EMU
        ],
        key=lambda s: s.left,
    )
    runs: List[List[BaseShape]] = []
    used = [False] * len(cands)
    for i, s in enumerate(cands):
        if used[i]:
            continue
        col = [s]
        used[i] = True
        for j in range(i + 1, len(cands)):
            if not used[j] and abs(cands[j].left - s.left) <= MARKER_X_TOL_EMU:
                col.append(cands[j])
                used[j] = True
        if len(col) >= MARKER_MIN_RUN:
            runs.append(sorted(col, key=lambda s: s.top))
    return runs


def _pair_run_to_text(
    run: List[BaseShape], text_shapes: List[BaseShape]
) -> Optional[BaseShape]:
    """The text box a marker run belongs to: the nearest text shape (within
    MARKER_PAIR_GAP_EMU horizontally) whose vertical span overlaps the run."""
    run_left = min(s.left for s in run)
    run_right = max(s.left + s.width for s in run)
    run_top = min(s.top for s in run)
    run_bot = max(s.top + s.height for s in run)
    best, best_gap = None, None
    for t in text_shapes:
        if min(run_bot, t.top + t.height) - max(run_top, t.top) <= 0:
            continue  # no vertical overlap
        if t.left >= run_right:
            gap = t.left - run_right
        elif t.left + t.width <= run_left:
            gap = run_left - (t.left + t.width)
        else:
            gap = 0
        if gap <= MARKER_PAIR_GAP_EMU and (best_gap is None or gap < best_gap):
            best, best_gap = t, gap
    return best


def _overlap_finding(
    a: BaseShape,
    b: BaseShape,
    ea: BoxEmu,
    eb: BoxEmu,
    pen_emu: int,
    kind: str,
    a_text: bool,
    b_text: bool,
) -> Dict[str, object]:
    """The structured record the QA digest tiers on. Captures which shape
    overflowed onto which ("over" spilled onto "hit"), how much of the engulfed
    box that covers ON the penetration axis - so the raw inches read back as a
    percentage ("0.40in" -> "100% of #142's height") - the axis, and whether both
    shapes carry text. Text-on-text is the loud, low-false-positive case."""
    ix = min(ea[0] + ea[2], eb[0] + eb[2]) - max(ea[0], eb[0])
    iy = min(ea[1] + ea[3], eb[1] + eb[3]) - max(ea[1], eb[1])
    axis = "horizontally" if ix < iy else "vertically"
    # For a spill, the overflower is the shape whose box grew past its declared
    # bounds; "hit" is the neighbour it landed on. A plain peer overlap has no
    # overflower, so keep document order (a is "over", b is "hit").
    over, over_t, hit, hit_box, hit_t = a, a_text, b, eb, b_text
    if kind == "spill" and ea == (a.left, a.top, a.width, a.height):
        over, over_t, hit, hit_box, hit_t = b, b_text, a, ea, a_text
    extent = hit_box[2] if axis == "horizontally" else hit_box[3]
    coverage = min(1.0, pen_emu / extent) if extent else 0.0
    return {
        "over": over.shape_id,
        "hit": hit.shape_id,
        "kind": kind,
        "axis": axis,
        "pen_in": pen_emu / EMU_PER_INCH,
        "coverage": coverage,
        "text_on_text": bool(over_t and hit_t),
    }


def _load_font(size: int):
    """A legible bold font at `size`px: prefer a TrueType bold, then a sized
    default, falling back to the bitmap default the sandbox always ships."""
    from PIL import ImageFont

    for loader in (
        lambda: ImageFont.truetype("DejaVuSans-Bold.ttf", size),
        lambda: ImageFont.load_default(size=size),  # Pillow >= 10.1
        ImageFont.load_default,
    ):
        try:
            return loader()
        except Exception:  # noqa: BLE001 - try the next fallback
            continue
    return None


def _annotate_boxes(
    image_path: Path,
    slide: Slide,
    slide_w_emu: int,
    slide_h_emu: int,
    effective_boxes: Optional[Dict[int, BoxEmu]] = None,
    draw_boxes: bool = True,
):
    """Overlay each top-level shape's bounding box on the slide image and
    compute pixel metrics. Box positions are read from the file (exact even
    though text is LibreOffice-rendered); the rest is measured on the rendered
    pixels (colors/positions are faithful). Draws box outlines, an `#id` label
    just OUTSIDE each box (detail goes to stdout, not onto the image), a tint on
    text boxes, and a red wash over peer-overlap regions.

    `draw_boxes=False` computes the same findings but draws only the collision
    wash. That is the QA default: the per-shape tint and outlines answer "which
    box is #id" at the cost of the question QA actually asks - is this copy
    legible on this background, is this margin even - because a 16%-alpha wash
    over every text box and a 3px saturated outline around it change exactly the
    pixels being judged.

    `effective_boxes` maps a shape_id to a box (EMU) grown to wrap copy that
    overflows its declared box; when present the overlay draws and tests that
    box, so an overflowing text box visibly wraps its text and a spill onto a
    neighbour is caught (a containment that appears only after a box grows is
    spillover, not a designed fg/bg overlay, so it is surfaced not suppressed).

    Returns (out_path, findings) where findings has keys: "overlaps" (a list of
    _overlap_finding dicts - which shape spilled onto which, the kind, axis,
    coverage and text-on-text flag) and "markers" [(id, off_in)] (decorative
    markers in a run with no rendered text row near them). None if the image is
    unreadable."""
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return None
    if not slide_w_emu or not slide_h_emu:
        return None
    try:
        base = Image.open(image_path).convert("RGBA")
    except (OSError, ValueError):
        return None
    width_px, height_px = base.size
    sample = base.convert("RGB")  # clean pixels for metrics (before overlay)
    overlay = Image.new("RGBA", base.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    font = _load_font(14)

    findings = {"overlaps": [], "markers": []}
    ppi = width_px / (slide_w_emu / EMU_PER_INCH)
    effective_boxes = effective_boxes or {}

    def declared(sh) -> BoxEmu:
        return (sh.left, sh.top, sh.width, sh.height)

    def box_of(sh) -> BoxEmu:
        return effective_boxes.get(sh.shape_id) or declared(sh)

    def to_px_box(b: BoxEmu):
        l, t, w, h = b
        return (
            l / slide_w_emu * width_px,
            t / slide_h_emu * height_px,
            (l + w) / slide_w_emu * width_px,
            (t + h) / slide_h_emu * height_px,
        )

    def to_px(sh):
        return to_px_box(box_of(sh))

    def has_text(sh):
        return sh.has_text_frame and any(
            (p.text or "").strip() for p in sh.text_frame.paragraphs
        )

    shapes = [
        s for s in slide.shapes
        if None not in (s.left, s.top, s.width, s.height)
    ]

    # Shade collision zones on the effective (text-extent) boxes: a peer overlap
    # as before, plus a containment that exists ONLY after a box grew to wrap
    # overflowing text - that text spilled onto a neighbour (a containment that
    # also holds at declared sizes is an intentional fg/bg overlay, suppressed).
    for i, a in enumerate(shapes):
        for b in shapes[i + 1:]:
            # Banners/backgrounds (full-bleed bands, backdrops) intentionally sit
            # under content that overlaps them at declared geometry - layering,
            # not a collision. But an overflowing box that spilled onto a
            # full-span shape (a title grown onto a full-width subtitle) is a
            # real overflow, so suppress only the declared-overlap layering case.
            if _suppress_as_layering(
                declared(a), declared(b), slide_w_emu, slide_h_emu
            ):
                continue
            ea, eb = box_of(a), box_of(b)
            flag, pen, kind = _render_collision(declared(a), declared(b), ea, eb)
            if not flag:
                continue
            ax0, ay0, ax1, ay1 = to_px_box(ea)
            bx0, by0, bx1, by1 = to_px_box(eb)
            ix0, iy0 = max(ax0, bx0), max(ay0, by0)
            ix1, iy1 = min(ax1, bx1), min(ay1, by1)
            if ix1 > ix0 and iy1 > iy0:
                # Wash only on the --boxes diagnostic view. On the plain QA
                # render it is a lie of omission: the pdf-word check downstream
                # clears most of these pairs (a designed overlay, two boxes whose
                # text sits in opposite halves), so a red block appears over
                # perfectly good copy with no finding to explain it, and the
                # render stops being a faithful picture of the slide.
                if draw_boxes:
                    draw.rectangle([ix0, iy0, ix1, iy1], fill=(255, 0, 0, 130))
                findings["overlaps"].append(
                    _overlap_finding(
                        a, b, ea, eb, pen, kind, has_text(a), has_text(b)
                    )
                )

    # Decorative-marker alignment: each marker in a run should have a rendered
    # text row near it in the text box it pairs with. Detect each text box's row
    # positions once, then check the runs against the box each pairs with.
    text_shapes = [s for s in shapes if has_text(s)]
    rows_by_id = {
        s.shape_id: _text_row_centers(sample, to_px(s)) for s in text_shapes
    }
    for run in _marker_runs(shapes):
        paired = _pair_run_to_text(run, text_shapes)
        if paired is None:
            continue  # standalone decoration, not row-markers
        rows = rows_by_id.get(paired.shape_id) or []
        if not rows:
            continue
        for m in run:
            my = (m.top + m.height / 2) / slide_h_emu * height_px
            nearest = min(abs(yc - my) for yc in rows)
            if nearest / ppi > MARKER_ALIGN_TOL_IN:
                findings["markers"].append((m.shape_id, nearest / ppi))

    if not draw_boxes:
        out_path = image_path.with_name(image_path.stem + "-qa.png")
        try:
            Image.alpha_composite(base, overlay).convert("RGB").save(out_path)
        except (OSError, ValueError):
            return None
        return out_path, findings

    for i, shape in enumerate(shapes):
        x0, y0, x1, y1 = to_px(shape)
        color = _BOX_PALETTE[i % len(_BOX_PALETTE)]
        if has_text(shape):
            draw.rectangle([x0, y0, x1, y1], fill=color + (40,))
        draw.rectangle([x0, y0, x1, y1], outline=color + (255,), width=3)

    # Labels last, each on a small filled chip so the #id stays legible over any
    # content or outline. Place it just above its box (just inside the top when
    # there is no room above), then nudge it down past any already-placed label
    # so stacked boxes don't smear their labels together.
    placed: List[Tuple[float, float, float, float]] = []
    for i, shape in enumerate(shapes):
        x0, y0, _, _ = to_px(shape)
        color = _BOX_PALETTE[i % len(_BOX_PALETTE)]
        text = f"#{shape.shape_id}"
        try:
            bl, bt, br, bb = draw.textbbox((0, 0), text, font=font)
            tw, th, off_l, off_t = br - bl, bb - bt, bl, bt
        except Exception:  # noqa: BLE001 - measuring is best-effort
            tw, th, off_l, off_t = 7 * len(text), 11, 0, 0
        pad = 2
        cw, ch = tw + 2 * pad, th + 2 * pad
        lx = max(0.0, x0)
        ly = y0 - ch - 1 if y0 - ch - 1 >= 0 else y0 + 1
        rect = (lx, ly, lx + cw, ly + ch)
        # O(n^2) over labels; n = shapes per slide (< ~40), so this is cheap.
        for _ in range(8):
            if not any(_rects_overlap(rect, p) for p in placed):
                break
            ly = rect[3] + 1
            rect = (lx, ly, lx + cw, ly + ch)
        placed.append(rect)
        draw.rectangle(list(rect), fill=color + (235,))
        draw.text(
            (lx + pad - off_l, ly + pad - off_t),
            text, fill=_contrast_text(color), font=font,
        )

    out_path = image_path.with_name(image_path.stem + "-boxes.png")
    try:
        Image.alpha_composite(base, overlay).convert("RGB").save(out_path)
    except (OSError, ValueError):
        return None
    return out_path, findings
