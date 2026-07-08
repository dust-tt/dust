"""Typography for pptx_inspect: placeholder-type detection, the layout/master/theme
default-resolution chain, and per-segment / per-paragraph text formatting.
(A segment is a span of same-formatted characters - python-pptx calls these
`runs`, so API access keeps `.runs`.)

Depends only on the shared ooxml/utils helpers (and python-pptx) - no
dependency on the other pptx_* modules - so audit and the CLI build on it.
"""
from __future__ import annotations

import zipfile
from typing import Dict, List, Optional, Tuple

import ooxml
from pptx.shapes.base import BaseShape
from utils import TEXT_PREVIEW_LIMIT, ellipsize, flatten_text

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# Glyphs commonly typed as manual bullets at the start of a paragraph. When
# they appear inside a placeholder whose layout already renders a bullet,
# the result is a doubled marker ("● • text"); see the bullet-glyph lint.
BULLET_PREFIXES = ("•", "·", "*", "-", "–")


def _starts_with_bullet_glyph(text: str) -> bool:
    if len(text) < 2:
        return False
    if text[0] not in BULLET_PREFIXES:
        return False
    return text[1].isspace()


def placeholder_type(shape: BaseShape) -> Optional[str]:
    if not getattr(shape, "is_placeholder", False):
        return None
    # `placeholder_format` raises ValueError on non-placeholder shapes,
    # which we already filtered out above.
    pf = shape.placeholder_format
    t = getattr(pf, "type", None)
    if t is None:
        return None
    name = getattr(t, "name", None)
    return name.lower() if name else None


def font_argb(segment) -> Optional[str]:
    font = getattr(segment, "font", None)
    if font is None:
        return None
    color = getattr(font, "color", None)
    if color is None:
        return None
    # color.rgb raises AttributeError when the color is theme-based
    # rather than explicit RGB; we treat both as "not set".
    try:
        rgb = color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    if rgb is None:
        return None
    return str(rgb).upper()


def segment_color_token(segment) -> Optional[str]:
    """A segment's explicit color as a short token, or None when it inherits.
    Explicit RGB -> "#RRGGBB"; theme/scheme color -> its scheme name (e.g.
    "dk2"). Surfacing the scheme token (not just RGB) matters because a box's
    paragraphs often carry a deliberate color pattern (white title, scheme-color
    body) that is invisible if only literal RGB is reported."""
    argb = font_argb(segment)
    if argb:
        return "#" + (argb[2:] if len(argb) == 8 else argb)
    font = getattr(segment, "font", None)
    color = getattr(font, "color", None) if font is not None else None
    if color is None or getattr(color, "type", None) is None:
        return None
    theme = getattr(color, "theme_color", None)
    name = getattr(theme, "name", None)
    return name.lower() if isinstance(name, str) else None


# ---------------------------------------------------------------------------
# Theme / placeholder default resolution.
#
# python-pptx does not surface the effective typography of an empty layout
# placeholder - that information lives in the layout's <a:lstStyle>, falling
# back to the master's matching placeholder, the master's titleStyle /
# bodyStyle / otherStyle, and finally the theme's major/minor font + color
# scheme. The helpers below walk that chain directly off the zip so we can
# print "this title placeholder defaults to Lexend 28pt bold #F8FAFC".
# ---------------------------------------------------------------------------


def _qp(local: str) -> str:
    return f"{{{P_NS}}}{local}"


def _read_clr_map(master_xml) -> Dict[str, str]:
    """Map placeholder color tokens (bg1, tx1, accent1...) to theme tokens
    (lt1/dk1/accent1...) via the master's <p:clrMap>."""
    if master_xml is None:
        return {}
    clr_map = master_xml.find(_qp("clrMap"))
    if clr_map is None:
        return {}
    return {k: v for k, v in clr_map.attrib.items() if not k.startswith("{")}


def _resolve_scheme_color(
    scheme_token: str,
    clr_map: Dict[str, str],
    theme_colors: Dict[str, str],
) -> Optional[str]:
    target = clr_map.get(scheme_token, scheme_token)
    hex_val = theme_colors.get(target)
    return hex_val or None


def _solid_fill_hex(
    elem,
    clr_map: Dict[str, str],
    theme_colors: Dict[str, str],
) -> Optional[str]:
    """Resolve <a:solidFill> on a defRPr-like element to a 6-hex color."""
    if elem is None:
        return None
    fill = elem.find(ooxml.qa("solidFill"))
    if fill is None:
        return None
    srgb = fill.find(ooxml.qa("srgbClr"))
    if srgb is not None:
        return srgb.attrib.get("val", "").upper() or None
    scheme = fill.find(ooxml.qa("schemeClr"))
    if scheme is not None:
        return _resolve_scheme_color(scheme.attrib.get("val", ""), clr_map, theme_colors)
    return None


def _ph_type_default_kind(ph_type: Optional[str]) -> str:
    """Pick which master *Style block applies when nothing else matches."""
    if ph_type in ("title", "ctrtitle", "ctr_title", "center_title"):
        return "title"
    if ph_type in ("body", "subtitle", "subTitle", "obj"):
        return "body"
    return "other"


def _find_ph_sp(parent, ph_idx: Optional[int], ph_type: Optional[str]):
    """Find a <p:sp> in `parent` whose <p:ph> matches the given idx/type."""
    if parent is None:
        return None
    for sp in parent.iter(_qp("sp")):
        ph = sp.find(f"{_qp('nvSpPr')}/{_qp('nvPr')}/{_qp('ph')}")
        if ph is None:
            continue
        sp_idx_raw = ph.attrib.get("idx")
        sp_idx = int(sp_idx_raw) if sp_idx_raw and sp_idx_raw.isdigit() else 0
        sp_type = (ph.attrib.get("type") or "").lower()
        if ph_idx is not None and sp_idx == ph_idx:
            return sp
        if ph_type and sp_type == ph_type:
            return sp
    return None


def _lvl1_defrpr(sp):
    """Return (lvl1pPr_element, defRPr_element) from the <a:lstStyle> of a
    placeholder shape. Either may be None."""
    if sp is None:
        return None, None
    lst = sp.find(f"{_qp('txBody')}/{ooxml.qa('lstStyle')}")
    if lst is None:
        return None, None
    lvl1 = lst.find(ooxml.qa("lvl1pPr"))
    if lvl1 is None:
        return None, None
    def_rpr = lvl1.find(ooxml.qa("defRPr"))
    return lvl1, def_rpr


def _master_style_defrpr(master_xml, kind: str):
    """Pull <a:lvl1pPr>/<a:defRPr> from <p:titleStyle> / <p:bodyStyle> /
    <p:otherStyle> on the master."""
    if master_xml is None:
        return None, None
    style_name = {
        "title": "titleStyle",
        "body": "bodyStyle",
        "other": "otherStyle",
    }[kind]
    style = master_xml.find(f"{_qp('txStyles')}/{_qp(style_name)}")
    if style is None:
        return None, None
    lvl1 = style.find(ooxml.qa("lvl1pPr"))
    if lvl1 is None:
        return None, None
    def_rpr = lvl1.find(ooxml.qa("defRPr"))
    return lvl1, def_rpr


def resolve_placeholder_defaults(
    layout_xml,
    master_xml,
    theme_xml,
    clr_map: Dict[str, str],
    theme_colors: Dict[str, str],
    ph_idx: Optional[int],
    ph_type: Optional[str],
) -> Dict[str, Optional[str]]:
    """Walk layout placeholder -> master placeholder -> master *Style ->
    theme to compute the effective default typography of a placeholder."""
    kind = _ph_type_default_kind(ph_type)

    chain = []
    chain.append(_lvl1_defrpr(_find_ph_sp(layout_xml, ph_idx, ph_type)))
    chain.append(_lvl1_defrpr(_find_ph_sp(master_xml, ph_idx, ph_type)))
    chain.append(_master_style_defrpr(master_xml, kind))

    typeface: Optional[str] = None
    size_pt: Optional[float] = None
    bold: Optional[bool] = None
    color_hex: Optional[str] = None
    align: Optional[str] = None

    for lvl1, def_rpr in chain:
        if lvl1 is not None and align is None:
            algn = lvl1.attrib.get("algn")
            if algn:
                align = algn
        if def_rpr is None:
            continue
        if typeface is None:
            latin = def_rpr.find(ooxml.qa("latin"))
            if latin is not None:
                typeface = latin.attrib.get("typeface") or None
        if size_pt is None:
            sz = def_rpr.attrib.get("sz")
            if sz and sz.isdigit():
                size_pt = int(sz) / 100.0
        if bold is None:
            b = def_rpr.attrib.get("b")
            if b in ("1", "true"):
                bold = True
            elif b in ("0", "false"):
                bold = False
        if color_hex is None:
            color_hex = _solid_fill_hex(def_rpr, clr_map, theme_colors)

    if typeface is None:
        typeface = ooxml.theme_font(
            theme_xml, "major" if kind == "title" else "minor")

    return {
        "typeface": typeface,
        "size_pt": size_pt,
        "bold": bold,
        "color": color_hex,
        "align": align,
    }


def format_placeholder_defaults(defaults: Dict[str, Optional[str]]) -> str:
    parts: List[str] = []
    if defaults.get("typeface"):
        parts.append(str(defaults["typeface"]))
    size_pt = defaults.get("size_pt")
    if size_pt is not None:
        if float(size_pt).is_integer():
            parts.append(f"{int(size_pt)}pt")
        else:
            parts.append(f"{size_pt:.1f}pt")
    if defaults.get("bold"):
        parts.append("bold")
    color = defaults.get("color")
    if color:
        parts.append(f"#{color}")
    align = defaults.get("align")
    if align:
        parts.append(f"algn={align}")
    return " ".join(parts)


# ---------------------------------------------------------------------------
# Package-level XML access (lazy: most code paths don't need it).
# ---------------------------------------------------------------------------


def _slide_master_rel_target(zf: zipfile.ZipFile, master_path: str, rel_type_tail: str) -> Optional[str]:
    rels = ooxml.parse_rels(zf, ooxml.rels_path_for(master_path))
    if not rels:
        return None
    rels_xml = ooxml.read_xml(zf, ooxml.rels_path_for(master_path))
    if rels_xml is None:
        return None
    for rel in rels_xml.findall("pr:Relationship", ooxml.NS):
        rtype = rel.attrib.get("Type", "")
        if rtype.endswith(rel_type_tail):
            return ooxml.resolve_rel_target(ooxml.rels_path_for(master_path), rel.attrib["Target"])
    return None


def _read_theme_for_master(zf: zipfile.ZipFile, master_path: str):
    theme_path = _slide_master_rel_target(zf, master_path, "/theme")
    if not theme_path:
        return None, None
    return theme_path, ooxml.read_xml(zf, theme_path)


def _layout_master_path(zf: zipfile.ZipFile, layout_path: str) -> Optional[str]:
    return _slide_master_rel_target(zf, layout_path, "/slideMaster")


def _read_layout_chain(
    zf: zipfile.ZipFile, layout_path: str
) -> Tuple[Optional[object], Optional[object], Optional[object], Dict[str, str], Dict[str, str]]:
    layout_xml = ooxml.read_xml(zf, layout_path)
    master_path = _layout_master_path(zf, layout_path)
    master_xml = ooxml.read_xml(zf, master_path) if master_path else None
    theme_xml = None
    if master_path:
        _, theme_xml = _read_theme_for_master(zf, master_path)
    clr_map = _read_clr_map(master_xml)
    theme_colors = ooxml.theme_colors_by_name(theme_xml)
    return layout_xml, master_xml, theme_xml, clr_map, theme_colors


def font_hints(segment) -> str:
    font = getattr(segment, "font", None)
    if font is None:
        return ""
    parts: List[str] = []
    name = getattr(font, "name", None)
    if name:
        parts.append(name)
    size = getattr(font, "size", None)
    if size is not None:
        parts.append(f"{int(size.pt)}pt")
    if getattr(font, "bold", None):
        parts.append("bold")
    if getattr(font, "italic", None):
        parts.append("italic")
    if getattr(font, "underline", None):
        parts.append("underline")
    token = segment_color_token(segment)
    if token and token != "#000000":
        parts.append(f"color:{token}")
    return ", ".join(parts)


def paragraph_segments_summary(paragraph) -> str:
    """First segment's font hints, as a stand-in for paragraph formatting."""
    for segment in paragraph.runs:
        hints = font_hints(segment)
        if hints:
            return hints
    return ""


def paragraph_alignment(paragraph) -> Optional[str]:
    alignment = getattr(paragraph, "alignment", None)
    if alignment is None:
        return None
    name = getattr(alignment, "name", None)
    return name.lower() if name else None


def text_frame_lines(shape: BaseShape, indent: str = "  ") -> List[str]:
    if not shape.has_text_frame:
        return []
    is_placeholder = placeholder_type(shape) is not None
    paras = list(shape.text_frame.paragraphs)
    # Paragraphs are addressed by index (text_frame.paragraphs[i]), so label them
    # p[i] - not by bullet level, which is not an address. Interleaved EMPTY
    # paragraphs (spacer lines between content rows) are part of the box's
    # skeleton: they dictate which slots new text must land in and which carry
    # which styling, so surface them rather than dropping them. Trailing empties
    # carry no skeleton signal, so collapse them to a single count.
    last_content = -1
    for i, p in enumerate(paras):
        if flatten_text(p.text or "").strip():
            last_content = i
    lines: List[str] = []
    for i, paragraph in enumerate(paras):
        if i > last_content:
            break
        text = flatten_text(paragraph.text or "").strip()
        level = paragraph.level or 0
        attrs: List[str] = []
        if level:
            attrs.append(f"L{level}")
        hints = paragraph_segments_summary(paragraph)
        if hints:
            attrs.append(hints)
        algn = paragraph_alignment(paragraph)
        if algn:
            attrs.append(f"algn={algn}")
        if text and is_placeholder and _starts_with_bullet_glyph(text):
            attrs.append("[!] manual bullet glyph in placeholder")
        body = ellipsize(text, TEXT_PREVIEW_LIMIT) if text else "(empty)"
        line = f"{indent}p[{i}]: {body}"
        if attrs:
            line += f"  [{', '.join(attrs)}]"
        lines.append(line)
    trailing = len(paras) - 1 - last_content
    if last_content >= 0 and trailing > 0:
        lines.append(f"{indent}p[{last_content + 1}..{len(paras) - 1}]: (empty x{trailing})")
    return lines
