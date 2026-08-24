"""Audit layer for pptx_inspect: per-shape lint markers (overflow, overlap,
distorted/low-res images, fit/overset, cover detection) and deck-fidelity
analysis against a template (imagery, bare-canvas, density, structural tally,
content-slot fidelity, shape retention).

Builds on geometry + typography; produces the [!]/[i] findings the CLI formats
and the inputs to the --compare [QA: PASS/FAIL] verdict.
"""
from __future__ import annotations

import re
import zipfile
from typing import Dict, Iterable, List, NamedTuple, Optional, Tuple

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.presentation import Presentation as PresentationType
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

from utils import flatten_text
from pptx_geometry import (
    ASPECT_TOLERANCE,
    CENTER_OFFSET_EMU,
    EDGE_EPSILON_EMU,
    EMU_PER_INCH,
    FILL_FLOATING,
    FULL_SPAN,
    MIN_IMAGE_DPI,
    OVERSET_TOLERANCE,
    SAFE_MARGIN_EMU,
    _classify_overlap,
    _fit_estimate,
    _frame_text_len,
    _grows_to_fit,
    emu_to_inches,
    shape_kind,
)
from pptx_typography import placeholder_type, resolve_placeholder_defaults


# Drawing-namespace handles for detecting embedded images anywhere in a shape
# subtree. A populated picture placeholder reports shape_type PLACEHOLDER (not
# PICTURE), and a picture *fill* on an auto shape or the slide background is not
# a picture shape at all - but all of them carry a populated <a:blip>.
A_BLIP = qn("a:blip")
R_EMBED = qn("r:embed")
R_LINK = qn("r:link")
P_BG = qn("p:bg")


def _has_embedded_blip(element) -> bool:
    """True if the element subtree contains a populated <a:blip> (an embedded or
    linked image). Catches imagery that shape_type alone misses."""
    if element is None:
        return False
    for blip in element.iter(A_BLIP):
        if blip.get(R_EMBED) or blip.get(R_LINK):
            return True
    return False


class SlideContext(NamedTuple):
    width_emu: int
    height_emu: int


CoverRect = Tuple[int, int, int, int, int]  # (shape_id, left, top, width, height)


def _cover_candidates(shapes: Iterable[BaseShape]) -> List[CoverRect]:
    """Bounding boxes of non-placeholder shapes on the slide, used to decide
    whether an empty placeholder is actually covered by visible content."""
    out: List[CoverRect] = []
    for shape in shapes:
        if placeholder_type(shape) is not None:
            continue
        left, top, width, height = (
            shape.left,
            shape.top,
            shape.width,
            shape.height,
        )
        if None in (left, top, width, height):
            continue
        out.append((shape.shape_id, left, top, width, height))
    return out


def _find_covering_shape(
    placeholder: BaseShape, candidates: List[CoverRect]
) -> Optional[int]:
    """Return the shape_id of a non-placeholder shape whose bounding box
    covers ≥50% of the placeholder's box, or None if no shape does. Used to
    flip the empty-placeholder marker from "populate" to "delete"."""
    pl_left = placeholder.left
    pl_top = placeholder.top
    pl_w = placeholder.width
    pl_h = placeholder.height
    if None in (pl_left, pl_top, pl_w, pl_h) or pl_w <= 0 or pl_h <= 0:
        return None
    pl_right = pl_left + pl_w
    pl_bottom = pl_top + pl_h
    pl_area = pl_w * pl_h
    for shape_id, left, top, w, h in candidates:
        ix = max(0, min(pl_right, left + w) - max(pl_left, left))
        iy = max(0, min(pl_bottom, top + h) - max(pl_top, top))
        if ix * iy * 2 >= pl_area:
            return shape_id
    return None


def _effective_typeface(shape: BaseShape, layout_chain) -> Optional[str]:
    """The face this shape's text renders in: an explicit run typeface if set,
    else the placeholder's resolved layout default."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                name = getattr(getattr(run, "font", None), "name", None)
                if name:
                    return name
    ph = placeholder_type(shape)
    if ph and layout_chain is not None:
        layout_xml, master_xml, theme_xml, clr_map, theme_colors = layout_chain
        pf = getattr(shape, "placeholder_format", None)
        idx = pf.idx if pf else None
        defaults = resolve_placeholder_defaults(
            layout_xml, master_xml, theme_xml, clr_map, theme_colors, idx, ph
        )
        face = defaults.get("typeface")
        if face:
            return str(face)
    return None


def _effective_font_size_pt(shape: BaseShape, layout_chain) -> Optional[float]:
    """The size text in this shape renders at: an explicit run size if set,
    else the placeholder's resolved layout default."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                size = getattr(run.font, "size", None)
                if size is not None:
                    return size.pt
    ph = placeholder_type(shape)
    if ph and layout_chain is not None:
        layout_xml, master_xml, theme_xml, clr_map, theme_colors = layout_chain
        pf = getattr(shape, "placeholder_format", None)
        idx = pf.idx if pf else None
        defaults = resolve_placeholder_defaults(
            layout_xml, master_xml, theme_xml, clr_map, theme_colors, idx, ph
        )
        size_pt = defaults.get("size_pt")
        if size_pt:
            return float(size_pt)
    return None


def _fit_tokens(shape: BaseShape, layout_chain) -> List[str]:
    """A capacity annotation for any sized text shape, plus an overset marker
    when text grossly exceeds a real container's estimated capacity."""
    used = _frame_text_len(shape)
    if used == 0:
        return []
    size_pt = _effective_font_size_pt(shape, layout_chain)
    if size_pt is None:
        return []
    est = _fit_estimate(shape, size_pt, _effective_typeface(shape, layout_chain))
    if est is None:
        return []
    size_label = (
        int(size_pt) if float(size_pt).is_integer() else round(size_pt, 1)
    )
    if est.capacity is None:
        return [f"~{est.chars_per_line}ch/line@{size_label}pt"]
    about = "" if est.measured else "~"
    tokens = [f"holds{about}{est.capacity}ch@{size_label}pt"]
    if not _grows_to_fit(shape):
        if used > est.capacity * OVERSET_TOLERANCE:
            tokens.append(
                f"[!] text overset (est) - ~{used} chars; box holds "
                f"~{est.capacity} at {size_label}pt"
            )
        elif used < est.capacity * FILL_FLOATING:
            tokens.append(
                f"[i] underfilled - ~{used} of ~{est.capacity}ch used"
            )
    return tokens


def embedded_image(shape: BaseShape):
    """A shape's image part, or None. python-pptx RAISES on `.image` for a
    picture placeholder with no picture in it, which `getattr(..., None)` does
    not catch - and a template full of empty picture placeholders is exactly
    what an audit walks."""
    try:
        return getattr(shape, "image", None)
    except (ValueError, KeyError, AttributeError):
        return None


def _image_markers(shape: BaseShape) -> List[str]:
    """Aspect-ratio distortion and low-resolution warnings for a picture (or a
    populated picture placeholder). Distortion compares the display box ratio to
    the image's crop-adjusted native ratio - squished/stretched photos read as
    sloppy even when nothing overflows."""
    image = embedded_image(shape)
    if image is None:
        return []
    size = getattr(image, "size", None)
    if not (isinstance(size, tuple) and len(size) == 2):
        return []
    nat_w, nat_h = size
    w, h = shape.width, shape.height
    if not w or not h or nat_w <= 0 or nat_h <= 0:
        return []
    try:
        cl = float(getattr(shape, "crop_left", 0) or 0)
        cr = float(getattr(shape, "crop_right", 0) or 0)
        ct = float(getattr(shape, "crop_top", 0) or 0)
        cb = float(getattr(shape, "crop_bottom", 0) or 0)
    except (TypeError, ValueError):
        cl = cr = ct = cb = 0.0
    vis_w = nat_w * max(1e-6, 1.0 - cl - cr)
    vis_h = nat_h * max(1e-6, 1.0 - ct - cb)
    nat_ratio = vis_w / vis_h
    disp_ratio = w / h
    markers: List[str] = []
    if nat_ratio > 0:
        rel = disp_ratio / nat_ratio
        if rel > ASPECT_TOLERANCE or rel < 1 / ASPECT_TOLERANCE:
            markers.append(
                f"[!] image distorted - box {disp_ratio:.2f}:1 vs image native "
                f"{nat_ratio:.2f}:1"
            )
    disp_w_in = w / EMU_PER_INCH
    if disp_w_in > 0:
        eff_dpi = vis_w / disp_w_in
        if eff_dpi < MIN_IMAGE_DPI:
            markers.append(
                f"[!] low-res image - ~{int(eff_dpi)} dpi at display size"
            )
    return markers


def _overlap_markers(
    shape: BaseShape, all_boxes: Optional[List[CoverRect]]
) -> List[str]:
    """Flag this shape's overlaps quantitatively (each pair once, on the lower
    shape id): stacked -> [!] blocker; partial peer overlap -> [i] advisory with
    the inches to separate and the axis; containment (fg/bg) -> suppressed."""
    if not all_boxes:
        return []
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    if None in (left, top, width, height) or width <= 0 or height <= 0:
        return []
    out: List[str] = []
    for oid, ol, ot, ow, oh in all_boxes:
        if oid <= shape.shape_id or None in (ol, ot, ow, oh):
            continue
        kind, pen, axis = _classify_overlap(
            (left, top, width, height), (ol, ot, ow, oh)
        )
        if kind == "stacked":
            out.append(f"[!] stacked with shape #{oid} - boxes coincide")
        elif kind == "peer":
            out.append(
                f"[i] overlaps shape #{oid} by {pen / EMU_PER_INCH:.2f}in {axis}"
            )
        elif kind == "contained":
            # the smaller box sits inside the larger (fg/bg): note if the inner
            # box is markedly off-centre (a number floating high in its card).
            dx = abs((left + width / 2) - (ol + ow / 2))
            dy = abs((top + height / 2) - (ot + oh / 2))
            if max(dx, dy) >= CENTER_OFFSET_EMU:
                off = "horizontally" if dx >= dy else "vertically"
                out.append(
                    f"[i] off-centre from shape #{oid} by "
                    f"{max(dx, dy) / EMU_PER_INCH:.2f}in {off}"
                )
    return out


def _shape_warning_markers(
    shape: BaseShape,
    ph: Optional[str],
    ctx: Optional[SlideContext],
    cover_candidates: Optional[List[CoverRect]] = None,
    all_boxes: Optional[List[CoverRect]] = None,
) -> List[str]:
    markers: List[str] = []
    if ctx is not None:
        left, top, width, height = (
            shape.left,
            shape.top,
            shape.width,
            shape.height,
        )
        if None not in (left, top, width, height):
            if width <= 0 or height <= 0:
                # A box collapsed on one axis renders nothing at all: the text
                # is in the file, invisible on the slide, and every other check
                # here reads it as present. Comes from arithmetic that produced
                # a zero width/height, so it is always the model's own bug.
                markers.append(
                    f"[!] zero-size box ({emu_to_inches(width):.1f}x"
                    f"{emu_to_inches(height):.1f}in) - its content does not render"
                )
            elif (
                left < -EDGE_EPSILON_EMU
                or top < -EDGE_EPSILON_EMU
                or left + width > ctx.width_emu + EDGE_EPSILON_EMU
                or top + height > ctx.height_emu + EDGE_EPSILON_EMU
            ):
                markers.append("[!] extends past slide edge")
            elif ctx.width_emu and ctx.height_emu and shape.has_text_frame and any(
                (p.text or "").strip() for p in shape.text_frame.paragraphs
            ):
                # Text crowding a slide edge (advisory). Skip full-span banners,
                # which legitimately bleed to the edge.
                if (
                    width < FULL_SPAN * ctx.width_emu
                    and height < FULL_SPAN * ctx.height_emu
                ):
                    gap = min(
                        left, top,
                        ctx.width_emu - (left + width),
                        ctx.height_emu - (top + height),
                    )
                    if 0 <= gap < SAFE_MARGIN_EMU:
                        markers.append(
                            f"[i] text {gap / EMU_PER_INCH:.2f}in from nearest "
                            "slide edge"
                        )
    markers.extend(_image_markers(shape))
    markers.extend(_overlap_markers(shape, all_boxes))
    if ph and shape.has_text_frame:
        has_text = any(
            (p.text or "").strip() for p in shape.text_frame.paragraphs
        )
        if not has_text:
            cover_id = (
                _find_covering_shape(shape, cover_candidates)
                if cover_candidates
                else None
            )
            if cover_id is not None:
                markers.append(
                    f"[!] empty placeholder, covered by shape #{cover_id}"
                )
            else:
                markers.append(
                    "[!] empty placeholder (renders the layout's prompt text)"
                )
    return markers


# Filler that must never reach a delivered deck. The lorem list is deliberately
# not just "lorem ipsum": a template's second and third filler paragraphs start
# at "Ut enim ad minim veniam" and "Duis aute irure dolor", so matching only the
# opening words missed four slides of untouched filler sitting in the columns
# beside real copy. These tokens are Latin fragments no business deck contains.
_LOREM_TOKENS = (
    "lorem ipsum", "dolor sit amet", "consectetur", "adipiscing", "eiusmod",
    "incididunt", "quis nostrud", "ullamco", "commodo consequat", "duis aute",
    "reprehenderit", "voluptate velit", "cillum dolore", "fugiat nulla",
    "excepteur", "occaecat", "cupidatat", "proident", "officia deserunt",
    "mollit anim", "est laborum",
)

_LEFTOVER_RE = re.compile(
    r"^[\[<].*[\]>]$|^x{3,}$|click to add|<[^>]+>|"
    + "|".join(re.escape(t) for t in _LOREM_TOKENS),
    re.IGNORECASE,
)


def _is_leftover_suspect(text: str) -> bool:
    t = text.strip()
    return bool(t) and bool(_LEFTOVER_RE.search(t))


def _package_names(path: str, prefix: str) -> Optional[set]:
    """Set of zip entry names under `prefix` (skipping directory entries),
    or None if the file isn't a readable package."""
    try:
        zf = zipfile.ZipFile(path)
    except (zipfile.BadZipFile, OSError):
        return None
    with zf:
        return {
            name
            for name in zf.namelist()
            if name.startswith(prefix) and not name.endswith("/")
        }


def _count_slides(path: str) -> int:
    names = _package_names(path, "ppt/slides/slide") or set()
    return sum(
        1 for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
    )


# ---------------------------------------------------------------------------
# Template-fidelity audit (--compare).
#
# `--compare` is the authoritative structural QA gate. Beyond "did you drop
# media/fonts" it answers "does this deck still respect the template, or did
# you rebuild it as a wall of text on the template's background?" - the single
# most common way template edits go wrong. Each threshold is tuned to flag the
# rebuilt-as-text failure mode while leaving a faithful reskin alone.
# ---------------------------------------------------------------------------

# Template counts as image-rich above this fraction of slides carrying imagery.
IMAGERY_RICH_TEMPLATE = 0.34
# Flag when the output keeps less than this fraction of the template's per-slide
# imagery rate (catches "stripped every picture", output rate 0).
IMAGERY_STRIP_RATIO = 0.5
# A "bare" slide has no imagery and no placeholders - content hand-drawn on an
# empty canvas. Flag when bare slides exceed this fraction of the output...
BARE_RATE = 0.34
# ...and exceed the template's own bare-slide rate by this margin, so a template
# that itself ships bare (shape-built) slides doesn't fail a faithful deck.
BARE_MARGIN = 0.25
# Per-slide content checks (imagery, bare-canvas) are skipped for tiny extracts:
# a deck reduced to this many slides or fewer is a deliberate excerpt, not a
# rebuild, and should not be judged for dropping the template's imagery.
MIN_SLIDES_FOR_CONTENT_AUDIT = 2
# Flag slides whose word count exceeds the template's max by more than this
# factor (a little headroom so a slide a few words over isn't nagged).
DENSITY_TOLERANCE = 1.15
# Shape retention vs the cloned exemplar. Trimming surplus elements is expected
# (legitimate clean-room decks bottom out around 0.79 kept), so the advisory
# only fires well below that: a slide kept < this fraction AND dropped at least
# SHAPE_DROP_MIN shapes has been gutted rather than adapted. Advisory, not a
# blocker - heavy adaptation is sometimes the right call.
SHAPE_RETENTION_FLOOR = 0.6
SHAPE_DROP_MIN = 2


# Share of the slide an image must cover to count as the slide's imagery. The
# deck's own footer logo is ~0.16% and a small 3-up icon ~1.1%, so this floor
# keeps the icons and drops the furniture. Without it a deck that stripped every
# photo still scored "imagery on 11/11 slides" - the template's logo, cloned onto
# each slide - and passed the gate that exists to catch exactly that.
IMAGERY_MIN_AREA = 0.005


def _imagery_in_shape(shape: BaseShape, slide_area_emu: int = 0) -> int:
    """Count pictures, charts, tables, and any image-bearing shape (a picture
    placeholder or picture-filled auto shape), recursing into groups. Pictures
    smaller than IMAGERY_MIN_AREA of the slide are furniture, not imagery."""
    kind = shape_kind(shape)
    if kind == "group":
        return sum(_imagery_in_shape(inner, slide_area_emu) for inner in shape.shapes)
    if kind not in ("pic", "chart", "table") and not _has_embedded_blip(
        getattr(shape, "_element", None)
    ):
        return 0
    if slide_area_emu and kind not in ("chart", "table"):
        width, height = shape.width, shape.height
        if None not in (width, height):
            if width * height < slide_area_emu * IMAGERY_MIN_AREA:
                return 0
    return 1


def _slide_bg_has_image(slide: Slide) -> bool:
    """True if the slide sets its own picture background - a full-bleed image
    that reads as imagery but is not a shape."""
    bg = slide._element.cSld.find(P_BG)
    return _has_embedded_blip(bg) if bg is not None else False


class DeckFidelity(NamedTuple):
    total: int  # slides in the slide list (what renders), not orphaned parts
    imagery_slides: int
    imagery_objs: int
    bare_slides: int
    layout_counts: Dict[str, int]
    word_counts: List[int]


def _deck_fidelity(prs: PresentationType) -> DeckFidelity:
    """Per-deck structure used by the audit. python-pptx walks only the slide
    list, so orphaned slide parts are correctly excluded here."""
    layout_counts: Dict[str, int] = {}
    word_counts: List[int] = []
    imagery_slides = 0
    imagery_objs = 0
    bare_slides = 0
    slide_area = (prs.slide_width or 0) * (prs.slide_height or 0)
    for slide in prs.slides:
        name = slide.slide_layout.name or "?"
        layout_counts[name] = layout_counts.get(name, 0) + 1
        n = sum(_imagery_in_shape(s, slide_area) for s in slide.shapes)
        imagery_objs += n
        has_imagery = n > 0 or _slide_bg_has_image(slide)
        if has_imagery:
            imagery_slides += 1
        has_placeholder = any(
            getattr(s, "is_placeholder", False) for s in slide.shapes
        )
        if not has_imagery and not has_placeholder:
            bare_slides += 1
        word_counts.append(slide_word_count(slide))
    return DeckFidelity(
        total=len(prs.slides),
        imagery_slides=imagery_slides,
        imagery_objs=imagery_objs,
        bare_slides=bare_slides,
        layout_counts=layout_counts,
        word_counts=word_counts,
    )


def _listed_slide_count(path: str) -> Optional[int]:
    """Slides referenced by presentation.xml's <p:sldIdLst> - what actually
    renders - read without python-pptx, so the orphan check still works when a
    deck fails to parse."""
    try:
        zf = zipfile.ZipFile(path)
    except (zipfile.BadZipFile, OSError):
        return None
    with zf:
        try:
            pres = zf.read("ppt/presentation.xml").decode("utf-8", "ignore")
        except KeyError:
            return None
    m = re.search(r"<[\w]*:?sldIdLst[\s>](.*?)</[\w]*:?sldIdLst>", pres, re.S)
    if not m:
        return None
    return len(re.findall(r"<[\w]*:?sldId[\s>/]", m.group(1)))


STRUCTURAL_KINDS = ("zero_size", "off_slide", "distorted", "stacked")

# Per-kind label for the --compare structural line and its blocker message.
STRUCTURAL_LABELS = {
    "zero_size": "zero-size",
    "off_slide": "off-slide",
    "distorted": "distorted-img",
    "stacked": "stacked",
}


def _slide_structural(slide: Slide, sw: int, sh: int) -> Dict[str, List[int]]:
    """Structural faults on one slide, per kind, as the shape ids carrying
    them."""
    found: Dict[str, List[int]] = {kind: [] for kind in STRUCTURAL_KINDS}
    shapes = [
        s for s in slide.shapes
        if None not in (s.left, s.top, s.width, s.height)
    ]
    for i, a in enumerate(shapes):
        if a.width <= 0 or a.height <= 0:
            found["zero_size"].append(a.shape_id)
        elif sw and sh and (
            a.left < -EDGE_EPSILON_EMU or a.top < -EDGE_EPSILON_EMU
            or a.left + a.width > sw + EDGE_EPSILON_EMU
            or a.top + a.height > sh + EDGE_EPSILON_EMU
        ):
            found["off_slide"].append(a.shape_id)
        if any("distorted" in m for m in _image_markers(a)):
            found["distorted"].append(a.shape_id)
        for b in shapes[i + 1:]:
            kind, _, _ = _classify_overlap(
                (a.left, a.top, a.width, a.height),
                (b.left, b.top, b.width, b.height),
            )
            if kind == "stacked":
                found["stacked"].append(a.shape_id)
    return found


def _deck_structural_audit(
    file_path: str, source_path: Optional[str]
) -> Optional[Dict[str, List[Tuple[int, int, bool]]]]:
    """Structural faults across the deck as (slide_no, shape_id, inherited).

    `inherited` means the exemplar this slide was cloned from carries the SAME
    fault on the same shape id at the same geometry - the template shipped it.
    Both halves are needed: geometry alone would excuse a photo the model
    swapped into an inherited box at the wrong aspect ratio, and the fault alone
    would excuse a shape the model dragged further off the canvas.

    Baselining per shape rather than per deck matters: a template with ten
    decorations bleeding off the canvas would otherwise buy a free pass for ten
    real defects."""
    try:
        prs = Presentation(file_path)
        src_prs = Presentation(source_path) if source_path else None
    except Exception:  # noqa: BLE001
        return None
    sw, sh = prs.slide_width or 0, prs.slide_height or 0
    mapping = _exemplar_map(prs, src_prs) if src_prs else {}
    src_slides = list(src_prs.slides) if src_prs else []
    src_w, src_h = (
        (src_prs.slide_width or 0, src_prs.slide_height or 0)
        if src_prs
        else (0, 0)
    )
    src_geometry: Dict[int, Dict[int, Tuple[int, int, int, int]]] = {}
    src_faults: Dict[int, Dict[str, set]] = {}
    for i, slide in enumerate(src_slides, start=1):
        src_geometry[i] = {
            s.shape_id: (s.left, s.top, s.width, s.height) for s in slide.shapes
        }
        src_faults[i] = {
            kind: set(ids)
            for kind, ids in _slide_structural(slide, src_w, src_h).items()
        }

    out: Dict[str, List[Tuple[int, int, bool]]] = {
        kind: [] for kind in STRUCTURAL_KINDS
    }
    for slide_no, slide in enumerate(prs.slides, start=1):
        src_no = mapping.get(slide_no, -1)
        exemplar = src_geometry.get(src_no, {})
        faults = src_faults.get(src_no, {})
        geometry = {
            s.shape_id: (s.left, s.top, s.width, s.height) for s in slide.shapes
        }
        found = _slide_structural(slide, sw, sh)
        for kind, ids in found.items():
            for sid in ids:
                inherited = (
                    sid in faults.get(kind, ())
                    and exemplar.get(sid) == geometry.get(sid)
                )
                out[kind].append((slide_no, sid, inherited))
    return out


def _content_slot_indices(shape: BaseShape) -> set:
    """Indices of the non-empty paragraphs in a text frame - the slots that
    actually carry content (vs the empty spacer paragraphs between them)."""
    if not shape.has_text_frame:
        return set()
    return {
        i
        for i, p in enumerate(shape.text_frame.paragraphs)
        if flatten_text(p.text or "").strip()
    }


def _filled_spacer_slots(out_shape: BaseShape, src_shape: BaseShape) -> set:
    """Paragraph slots the output filled that are INTERIOR spacers in the
    template - empty in the template and sitting between its content slots.
    Writing here is the wrong-slot fill that breaks color (a new run inherits a
    darker default) and strands the markers pinned to the content rows. Returns
    empty unless the template box has a real skeleton (>=2 content slots), so a
    plain placeholder or a box filled with fewer items than the template never
    trips it."""
    src_content = _content_slot_indices(src_shape)
    if len(src_content) < 2:
        return set()
    lo, hi = min(src_content), max(src_content)
    return {
        i
        for i in _content_slot_indices(out_shape)
        if i not in src_content and lo < i < hi
    }


def _slot_audit(
    file_path: str, source_path: str
) -> List[Tuple[int, int, set, set]]:
    """Per-slide, per-shape content-slot comparison against the template. Each
    output slide is matched to the template slide it was cloned from by shape-id
    overlap (ids survive cloning); shared text boxes are then checked for fills
    that landed in the template's spacer paragraphs. Returns
    (slide_no, shape_id, filled_spacers, template_content_slots)."""
    try:
        out_prs = Presentation(file_path)
        src_prs = Presentation(source_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    src_slides = [
        ({sh.shape_id for sh in s.shapes},
         {sh.shape_id: sh for sh in s.shapes if sh.has_text_frame})
        for s in src_prs.slides
    ]
    findings: List[Tuple[int, int, set, set]] = []
    for out_no, out_slide in enumerate(out_prs.slides, start=1):
        out_ids = {sh.shape_id for sh in out_slide.shapes}
        best_map, best_overlap = None, 0
        for ids, text_map in src_slides:
            overlap = len(out_ids & ids)
            if overlap > best_overlap:
                best_overlap, best_map = overlap, text_map
        # <2 shared ids => no template counterpart (e.g. a from-scratch slide);
        # nothing to diff against.
        if best_map is None or best_overlap < 2:
            continue
        for sh in out_slide.shapes:
            src_sh = best_map.get(sh.shape_id) if sh.has_text_frame else None
            if src_sh is None:
                continue
            spacers = _filled_spacer_slots(sh, src_sh)
            if spacers:
                findings.append(
                    (out_no, sh.shape_id, spacers,
                     _content_slot_indices(src_sh))
                )
    return findings


def _drop_audit(
    file_path: str, source_path: str
) -> List[Tuple[int, int, int, int]]:
    """Per-slide shape retention vs the exemplar each output slide was cloned
    from (matched by shape-id overlap). Trimming surplus elements is expected;
    dropping most of an exemplar's shapes means the wrong exemplar was chosen.
    Returns (slide_no, src_slide_no, kept, dropped) for slides that dropped
    >= SHAPE_DROP_MIN shapes AND kept < SHAPE_RETENTION_FLOOR of the exemplar."""
    try:
        out_prs = Presentation(file_path)
        src_prs = Presentation(source_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    src_slides = [{sh.shape_id for sh in s.shapes} for s in src_prs.slides]
    findings: List[Tuple[int, int, int, int]] = []
    for out_no, out_slide in enumerate(out_prs.slides, start=1):
        out_ids = {sh.shape_id for sh in out_slide.shapes}
        best_i, best_overlap = None, 0
        for i, ids in enumerate(src_slides):
            overlap = len(out_ids & ids)
            if overlap > best_overlap:
                best_overlap, best_i = overlap, i
        if best_i is None or best_overlap < 2:
            continue
        src_ids = src_slides[best_i]
        kept = len(out_ids & src_ids)
        dropped = len(src_ids - out_ids)
        if dropped >= SHAPE_DROP_MIN and kept / len(src_ids) < SHAPE_RETENTION_FLOOR:
            findings.append((out_no, best_i + 1, kept, dropped))
    return findings


# Text repeated on this many output slides is the template's furniture (footer,
# confidentiality stamp, brand line), kept on purpose - never a leftover.
# Real furniture - a footer, a page number, a confidentiality stamp - is defined
# by WHERE it sits, not by how often it repeats: it comes from the master or the
# layout, or it is a small box hugging a slide edge. Counting repetitions instead
# was wrong in both directions: it excused filler that repeated because nobody
# rewrote it, and it excused a template's own body paragraph left on ten slides.
FURNITURE_EDGE_BAND_IN = 0.8
FURNITURE_MAX_AREA = 0.06
# A cloned slide with at least this many unchanged text shapes, making up at
# least this share of its text, was never authored - it is the template's slide
# with a new title on it (a style-guide page, a sample chart page). Individually
# each survivor is a judgment call; the whole slide is not.
UNTOUCHED_MIN_SHAPES = 3
UNTOUCHED_SHARE = 0.7
# Body copy repeated on this many slides is a broken deck whatever its origin.
REPEATED_TEXT_SLIDES = 3
# Short lines legitimately recur (a stage label, a column header); a sentence
# this long appearing verbatim on three slides does not.
REPEATED_TEXT_MIN_WORDS = 6
# How many leftover shapes to name before summarising the rest.
LEFTOVER_LISTED = 8


def _exemplar_map(out_prs, src_prs) -> Dict[int, int]:
    """Output slide number -> the 1-based template slide it was cloned from,
    matched on shape-id overlap (pptx_slides copies ids when it duplicates)."""
    src_ids = [{sh.shape_id for sh in s.shapes} for s in src_prs.slides]
    mapping: Dict[int, int] = {}
    for out_no, out_slide in enumerate(out_prs.slides, start=1):
        ids = {sh.shape_id for sh in out_slide.shapes}
        best_i, best_overlap = None, 0
        for i, candidate in enumerate(src_ids):
            overlap = len(ids & candidate)
            if overlap > best_overlap:
                best_overlap, best_i = overlap, i
        if best_i is not None and best_overlap >= 2:
            mapping[out_no] = best_i + 1
    return mapping


def _shape_texts(slide: Slide) -> Dict[int, str]:
    return {
        sh.shape_id: flatten_text(" ".join(_shape_text_iter(sh))).strip()
        for sh in slide.shapes
    }


def _chrome_text(prs: PresentationType) -> set:
    """Every text string the deck's layouts and masters carry. Genuine furniture
    is authored there, which is what separates it from a slide's own copy."""
    out: set = set()
    for master in prs.slide_masters:
        for holder in [master] + list(master.slide_layouts):
            for shape in holder.shapes:
                units = list(_shape_text_iter(shape))
                joined = flatten_text(" ".join(units)).strip()
                if joined:
                    out.add(joined)
                for raw in units:
                    text = flatten_text(raw).strip()
                    if text:
                        out.add(text)
    return out


def _is_furniture(shape: BaseShape, slide_w: int, slide_h: int) -> bool:
    """A small box hugging a slide edge: a footer, a page number, a stamp."""
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    if None in (left, top, width, height) or width <= 0 or height <= 0:
        return False
    if not slide_w or not slide_h:
        return False
    if width * height > slide_w * slide_h * FURNITURE_MAX_AREA:
        return False
    band = int(EMU_PER_INCH * FURNITURE_EDGE_BAND_IN)
    return (
        top <= band
        or left <= band
        or slide_h - (top + height) <= band
        or slide_w - (left + width) <= band
    )


def _repeated_text_audit(file_path: str) -> List[Tuple[str, List[int]]]:
    """Paragraphs that appear verbatim on several slides: (text, slide numbers).

    Per PARAGRAPH, not per shape: the failure this catches is a template
    paragraph surviving inside a box whose other paragraphs were rewritten, so
    every slide's shape text is unique while one line inside it repeats ten
    times. Ten slides carrying the same sentence is a broken deck whether it came
    from the template or from padding a deck to satisfy a gate.
    """
    try:
        prs = Presentation(file_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    slide_w, slide_h = prs.slide_width or 0, prs.slide_height or 0
    chrome = set()
    for text in _chrome_text(prs):
        chrome.add(text)
    where: Dict[str, List[int]] = {}
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if _is_furniture(shape, slide_w, slide_h):
                continue
            for raw in _shape_text_iter(shape):
                text = flatten_text(raw).strip()
                if len(text.split()) < REPEATED_TEXT_MIN_WORDS:
                    continue
                if text in chrome:
                    continue
                slides = where.setdefault(text, [])
                if slide_no not in slides:
                    slides.append(slide_no)
    return sorted(
        ((text, slides) for text, slides in where.items()
         if len(slides) >= REPEATED_TEXT_SLIDES),
        key=lambda item: -len(item[1]),
    )


def _leftover_copy_audit(
    file_path: str, source_path: str
) -> List[Tuple[int, int, str]]:
    """Shapes on a cloned slide still carrying the exemplar's own copy.

    A template's building-block slide ships scaffolding text - step numbers, a
    stage name, a section word, lorem - and an exemplar cloned for its layout
    keeps every shape the model did not explicitly rewrite. Those survivors
    render as confident-looking nonsense next to the real content, and nothing
    else in the audit sees them because the shape is neither empty nor
    overflowing. Text that repeats across several output slides is the
    template's furniture and is excluded.

    Returns (slide_no, shape_id, text)."""
    try:
        out_prs = Presentation(file_path)
        src_prs = Presentation(source_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    mapping = _exemplar_map(out_prs, src_prs)
    if not mapping:
        return []
    out_slides = list(out_prs.slides)
    src_texts = [_shape_texts(s) for s in src_prs.slides]
    chrome = _chrome_text(out_prs)
    slide_w, slide_h = out_prs.slide_width or 0, out_prs.slide_height or 0

    findings: List[Tuple[int, int, str]] = []
    for out_no, src_no in sorted(mapping.items()):
        exemplar = src_texts[src_no - 1]
        slide = out_slides[out_no - 1]
        by_id = {sh.shape_id: sh for sh in slide.shapes}
        for sid, text in _shape_texts(slide).items():
            if len(text) < 2 or text in chrome:
                continue
            shape = by_id.get(sid)
            if shape is not None and _is_furniture(shape, slide_w, slide_h):
                continue
            if exemplar.get(sid) == text:
                findings.append((out_no, sid, text))
    return findings


def untouched_slides(
    findings: List[Tuple[int, int, str]], file_path: str
) -> List[Tuple[int, int, int]]:
    """Slides whose text is mostly still the exemplar's: (slide_no, kept, total).

    Reads the leftover findings back against how much text the slide carries at
    all, so "three survivors" on a dense slide stays advisory while "three of
    three" is called what it is: a template slide that was cloned and shipped."""
    try:
        prs = Presentation(file_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    slide_w, slide_h = prs.slide_width or 0, prs.slide_height or 0
    chrome = _chrome_text(prs)
    per_slide: Dict[int, int] = {}
    for slide_no, slide in enumerate(prs.slides, start=1):
        total = 0
        for shape in slide.shapes:
            text = flatten_text(" ".join(_shape_text_iter(shape))).strip()
            if not text or text in chrome:
                continue
            if _is_furniture(shape, slide_w, slide_h):
                continue
            total += 1
        per_slide[slide_no] = total

    kept: Dict[int, int] = {}
    for slide_no, _sid, _text in findings:
        kept[slide_no] = kept.get(slide_no, 0) + 1

    out: List[Tuple[int, int, int]] = []
    for slide_no, n in sorted(kept.items()):
        total = per_slide.get(slide_no, 0)
        if n >= UNTOUCHED_MIN_SHAPES and total and n >= total * UNTOUCHED_SHARE:
            out.append((slide_no, n, total))
    return out


# A cloned slide covering less than this share of its exemplar's footprint has
# had its content removed rather than replaced, and renders as a title over a
# void. Shape counting misses it: dropping one full-bleed photo out of four
# shapes keeps 75% of the shapes and loses 60% of the slide.
COVERAGE_FLOOR = 0.5
# Below this absolute drop the difference is trimming, not a hole.
COVERAGE_MIN_DROP = 0.15


def _slide_coverage(slide: Slide, slide_w: int, slide_h: int) -> float:
    """Share of the canvas covered by the slide's shapes, by scanline union so
    overlapping boxes are not double-counted."""
    if not slide_w or not slide_h:
        return 0.0
    boxes = [
        (s.left, s.top, s.left + s.width, s.top + s.height)
        for s in slide.shapes
        if None not in (s.left, s.top, s.width, s.height)
        and s.width > 0
        and s.height > 0
    ]
    if not boxes:
        return 0.0
    edges = sorted({y for box in boxes for y in (box[1], box[3])})
    covered = 0
    for top, bottom in zip(edges, edges[1:]):
        spans = sorted(
            (box[0], box[2]) for box in boxes if box[1] <= top and box[3] >= bottom
        )
        merged_x = 0
        reach = None
        for left, right in spans:
            if reach is None or left > reach:
                merged_x += right - left
                reach = right
            elif right > reach:
                merged_x += right - reach
                reach = right
        covered += merged_x * (bottom - top)
    return covered / (slide_w * slide_h)


def _hole_audit(
    file_path: str, source_path: str
) -> List[Tuple[int, int, float, float]]:
    """Cloned slides that ended up mostly empty canvas: (slide_no, src_no,
    coverage, exemplar coverage)."""
    try:
        out_prs = Presentation(file_path)
        src_prs = Presentation(source_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    mapping = _exemplar_map(out_prs, src_prs)
    if not mapping:
        return []
    out_slides = list(out_prs.slides)
    src_slides = list(src_prs.slides)
    ow, oh = out_prs.slide_width or 0, out_prs.slide_height or 0
    sw, sh = src_prs.slide_width or 0, src_prs.slide_height or 0
    findings: List[Tuple[int, int, float, float]] = []
    for out_no, src_no in sorted(mapping.items()):
        out_cov = _slide_coverage(out_slides[out_no - 1], ow, oh)
        src_cov = _slide_coverage(src_slides[src_no - 1], sw, sh)
        if src_cov <= 0:
            continue
        if (
            out_cov < src_cov * COVERAGE_FLOOR
            and src_cov - out_cov >= COVERAGE_MIN_DROP
        ):
            findings.append((out_no, src_no, out_cov, src_cov))
    return findings


def _filler_audit(file_path: str) -> List[Tuple[int, int, str]]:
    """Every shape still holding template filler: lorem, a bracketed prompt, an
    unreplaced tag. Needs no template to compare against, and no judgment -
    filler in a delivered deck is always wrong."""
    try:
        prs = Presentation(file_path)
    except Exception:  # noqa: BLE001 - degrade visibly in the caller
        return []
    out: List[Tuple[int, int, str]] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            text = flatten_text(" ".join(_shape_text_iter(shape))).strip()
            if text and _is_leftover_suspect(text):
                out.append((slide_no, shape.shape_id, text))
    return out


def _shape_text_iter(shape: BaseShape) -> Iterable[str]:
    kind = shape_kind(shape)
    if kind == "group":
        for inner in shape.shapes:
            yield from _shape_text_iter(inner)
        return
    if kind == "table":
        for row in shape.table.rows:
            for cell in row.cells:
                yield cell.text or ""
        return
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            yield paragraph.text or ""


def slide_word_count(slide: Slide) -> int:
    count = 0
    for shape in slide.shapes:
        for text in _shape_text_iter(shape):
            count += len(text.split())
    return count
