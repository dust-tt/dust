"""Tier-2 tests for pptx_typography on a tiny programmatic deck (no checked-in
binary): the paragraph skeleton (text_frame_lines), segment-colour tokens, and
placeholder-type detection. The deep layout/master/theme resolution chain is
covered by the golden --slide / --layouts views against the real template.

Run directly (`python test_typography.py`) or under pytest.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches  # noqa: E402

import pptx_typography as T  # noqa: E402


def _textbox_with_skeleton():
    """A box whose paragraphs are: content, spacer, content, spacer, spacer -
    the interleaved-spacer + trailing-empties shape the skeleton must surface."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    tf = tb.text_frame
    tf.text = "Title"            # p0 content
    tf.add_paragraph()           # p1 spacer (interleaved)
    tf.add_paragraph().text = "Body"  # p2 content
    tf.add_paragraph()           # p3 spacer (trailing)
    tf.add_paragraph()           # p4 spacer (trailing)
    return tb


def test_skeleton_indices_and_spacers():
    lines = [ln.strip() for ln in T.text_frame_lines(_textbox_with_skeleton())]
    # interleaved spacer shown; the two trailing empties collapse to a count
    assert any(ln.startswith("p[0]: Title") for ln in lines)
    assert any(ln.startswith("p[1]: (empty)") for ln in lines)
    assert any(ln.startswith("p[2]: Body") for ln in lines)
    assert any("p[3..4]: (empty x2)" in ln for ln in lines)
    assert len(lines) == 4  # p0, p1, p2, trailing-collapse


def test_segment_color_token_explicit_rgb():
    tb = _textbox_with_skeleton()
    segment = tb.text_frame.paragraphs[0].runs[0]
    segment.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    assert T.segment_color_token(segment) == "#FFFFFF"


def test_segment_color_token_inherited_is_none():
    tb = _textbox_with_skeleton()
    segment = tb.text_frame.paragraphs[0].runs[0]
    assert T.segment_color_token(segment) is None  # no explicit colour set


def test_segment_color_token_black_is_reported():
    # segment_color_token reports explicit black; font_hints is what suppresses it
    tb = _textbox_with_skeleton()
    segment = tb.text_frame.paragraphs[0].runs[0]
    segment.font.color.rgb = RGBColor(0, 0, 0)
    assert T.segment_color_token(segment) == "#000000"


def test_placeholder_type_none_for_textbox():
    assert T.placeholder_type(_textbox_with_skeleton()) is None


if __name__ == "__main__":
    tests = [v for k, v in sorted(globals().items())
             if k.startswith("test_") and callable(v)]
    for fn in tests:
        fn()
        print(f"ok   {fn.__name__}")
    print(f"\n{len(tests)} typography tests passed")
