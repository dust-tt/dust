"""Tier-2 tests for pptx_validate: each check fires on a deck broken the way an
edit script breaks it, and a clean deck stays clean. Run directly
(`python test_validate.py`) or under pytest.

Lives in soffice/tests/, a subdir getLocalDirContent skips: it copies only the
regular files directly in soffice/ (never recursing), so tests never ship in the image. It adds soffice/ to sys.path to import the module.
"""
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402

import pptx_validate as V  # noqa: E402


def _deck(path: Path, slides: int = 3) -> None:
    prs = Presentation()
    for _ in range(slides):
        prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(path)


def _rewrite(src: Path, dest: Path, mutate) -> None:
    """Copy the package, letting `mutate(name, data)` drop or edit each part."""
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(
        dest, "w", zipfile.ZIP_DEFLATED
    ) as zout:
        for info in zin.infolist():
            out = mutate(info.filename, zin.read(info.filename))
            if out is not None:
                zout.writestr(info, out)


def test_clean_deck_is_valid():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck)
        assert V.check(str(deck)) == []
        assert "[VALID]" in V.report(str(deck))


def test_orphan_slide_part_is_reported():
    with tempfile.TemporaryDirectory() as d:
        deck, broken = Path(d) / "deck.pptx", Path(d) / "broken.pptx"
        _deck(deck)
        # Drop slide 2 from the slide order only, the way a hand-edited
        # sldIdLst does: the part stays behind as an orphan.
        def mutate(name, data):
            if name != "ppt/presentation.xml":
                return data
            xml = data.decode()
            start = xml.index("<p:sldId ")
            end = xml.index("/>", start) + 2
            return (xml[:start] + xml[end:]).encode()

        _rewrite(deck, broken, mutate)
        problems = V.check(str(broken))
        assert any("not in the deck" in p for p in problems), problems
        assert "[INVALID" in V.report(str(broken))


def test_missing_relationship_target_is_reported():
    with tempfile.TemporaryDirectory() as d:
        deck, broken = Path(d) / "deck.pptx", Path(d) / "broken.pptx"
        _deck(deck)
        _rewrite(
            deck, broken,
            lambda name, data: None if name == "ppt/slides/slide2.xml" else data,
        )
        assert any("points at missing" in p for p in V.check(str(broken)))


def test_dangling_rid_in_a_part_is_reported():
    with tempfile.TemporaryDirectory() as d:
        deck, broken = Path(d) / "deck.pptx", Path(d) / "broken.pptx"
        _deck(deck)

        def mutate(name, data):
            if name != "ppt/_rels/presentation.xml.rels":
                return data
            xml = data.decode()
            start = xml.index("<Relationship ")
            end = xml.index("/>", start) + 2
            return (xml[:start] + xml[end:]).encode()

        _rewrite(deck, broken, mutate)
        assert any("has no relationship" in p for p in V.check(str(broken)))


def test_malformed_part_is_reported():
    with tempfile.TemporaryDirectory() as d:
        deck, broken = Path(d) / "deck.pptx", Path(d) / "broken.pptx"
        _deck(deck)
        _rewrite(
            deck, broken,
            lambda name, data: b"<p:sld><unclosed>"
            if name == "ppt/slides/slide1.xml" else data,
        )
        assert any("not well-formed" in p for p in V.check(str(broken)))


def test_duplicate_shape_id_is_reported():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck, slides=1)
        prs = Presentation(deck)
        slide = prs.slides[0]
        from pptx.util import Inches
        a = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        b = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1))
        for el in b._element.iter():
            if el.tag.endswith("}cNvPr"):
                el.set("id", a._element.find(
                    ".//{http://schemas.openxmlformats.org/presentationml/"
                    "2006/main}cNvPr").get("id"))
                break
        prs.save(deck)
        assert any("duplicate shape id" in p for p in V.check(str(deck)))


def test_template_problems_are_not_reported_as_yours():
    with tempfile.TemporaryDirectory() as d:
        template, out = Path(d) / "t.pptx", Path(d) / "out.pptx"
        _deck(template)
        broken = Path(d) / "broken.pptx"
        _rewrite(
            template, broken,
            lambda name, data: b"<p:sld><unclosed>"
            if name == "ppt/slides/slide1.xml" else data,
        )
        shutil.copyfile(broken, out)
        problems, inherited = V.check_against(str(out), str(broken))
        assert problems == []
        assert inherited > 0


def test_unreadable_file_is_reported():
    with tempfile.TemporaryDirectory() as d:
        junk = Path(d) / "junk.pptx"
        junk.write_bytes(b"not a zip")
        assert V.check(str(junk)) == ["not a readable .pptx package: File is not a zip file"]


if __name__ == "__main__":
    tests = [v for k, v in sorted(globals().items())
             if k.startswith("test_") and callable(v)]
    for fn in tests:
        fn()
        print(f"ok   {fn.__name__}")
    print(f"\n{len(tests)} validate tests passed")
