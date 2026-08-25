"""Tier-2 tests for pptx_audit's signature deterministic logic: the content-slot
detector that powers the --compare spacer-fill blocker. Built on tiny
programmatic decks. The estimated lints (overlap/fit/image) live in geometry and
are covered there; the deck-fidelity gate is integration-tested by the golden
--compare view.

Run directly (`python test_audit.py`) or under pytest.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

import pptx_audit as A  # noqa: E402


def box(content_idx, n_paras):
    """A textbox of n paragraphs with non-empty text only at content_idx."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    tf = tb.text_frame
    tf.text = "x" if 0 in content_idx else ""
    for i in range(1, n_paras):
        p = tf.add_paragraph()
        if i in content_idx:
            p.text = "x"
    return tb


def test_content_slot_indices():
    assert A._content_slot_indices(box({0, 2}, 4)) == {0, 2}


def test_spacer_write_detected():
    # template skeleton: content on even slots, spacers on odd; output filled
    # consecutively -> p1, p3 land in interior spacer slots
    src = box({0, 2, 4, 6, 8}, 10)
    out = box({0, 1, 2, 3}, 10)
    assert A._filled_spacer_slots(out, src) == {1, 3}


def test_no_spacer_write_when_fewer_items():
    # filling a subset of the template's own content slots is legitimate
    src = box({0, 2, 4, 6, 8}, 10)
    out = box({0, 2, 4}, 10)
    assert A._filled_spacer_slots(out, src) == set()


def test_no_spacer_write_for_simple_placeholder():
    # a template box with <2 content slots has no interior skeleton to violate
    src = box({0}, 3)
    out = box({0, 1}, 3)
    assert A._filled_spacer_slots(out, src) == set()


def test_spacer_write_only_counts_interior():
    # a slot past the template's last content row is not an interior spacer
    src = box({0, 2}, 6)        # content slots {0,2}, lo=0 hi=2
    out = box({0, 2, 4}, 6)     # p4 is beyond hi -> not flagged here
    assert A._filled_spacer_slots(out, src) == set()




def test_filler_detector_catches_lorem_past_its_opening_words():
    """A template's second and third filler paragraphs never say "lorem ipsum";
    matching only the opening words let four slides of untouched filler ship."""
    assert A._is_leftover_suspect("Lorem ipsum dolor sit amet, consectetur.")
    assert A._is_leftover_suspect(
        "Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris."
    )
    assert A._is_leftover_suspect(
        "Duis aute irure dolor in reprehenderit in voluptate velit esse."
    )
    assert A._is_leftover_suspect("[Slide Title]")
    assert A._is_leftover_suspect("Click to add text")


def test_filler_detector_leaves_real_copy_alone():
    for line in (
        "70+ connectors and MCP servers",
        "Start with one workflow. Expand from there.",
        "Pods and Frames for shared work",
        "EUR 24 per seat per month",
        "Contrat de vente, douleur et suivi",
    ):
        assert not A._is_leftover_suspect(line), line

def test_furniture_is_a_small_box_at_a_slide_edge_not_repeated_text():
    """Repetition used to define furniture, which excused a template body
    paragraph left on ten slides. Position defines it instead."""
    class Box:
        def __init__(self, left, top, width, height):
            self.left, self.top = left, top
            self.width, self.height = width, height

    inch = 914400
    slide_w, slide_h = int(13.33 * inch), int(7.5 * inch)
    footer = Box(int(0.5 * inch), int(7.0 * inch), int(3 * inch), int(0.3 * inch))
    body = Box(int(0.29 * inch), int(1.32 * inch), int(8.98 * inch), int(0.95 * inch))
    assert A._is_furniture(footer, slide_w, slide_h)
    assert not A._is_furniture(body, slide_w, slide_h)


def test_repeated_text_needs_a_real_sentence():
    assert A.REPEATED_TEXT_MIN_WORDS >= 6, (
        "short recurring labels (a stage name, a column header) are legitimate"
    )
    assert A.REPEATED_TEXT_SLIDES >= 3


def test_padding_is_baselined_against_the_template_not_a_bare_count():
    """Three slides sharing one picture is the padding shape we actually see (a
    case-study logo and its chart pasted onto the next three slides), so the
    count alone cannot be the rule - the template's own usage has to excuse an
    icon set it repeats itself."""
    assert A.REPEATED_IMAGE_SLIDES == 3


def test_imagery_min_area_excludes_a_footer_logo_and_keeps_an_icon():
    inch = 914400
    slide_area = int(10 * inch) * int(5.625 * inch)
    logo = (int(0.3 * inch)) ** 2
    icon = (int(0.8 * inch)) ** 2
    assert logo < slide_area * A.IMAGERY_MIN_AREA
    assert icon > slide_area * A.IMAGERY_MIN_AREA


def test_leading_break_detector_matches_every_break_character():
    """python-pptx surfaces <a:br/> as a vertical tab, which is what the model
    actually writes; newline and carriage return are covered for the same
    reason."""
    for ch in ("\v", "\n", "\r"):
        assert ch in A._LINE_BREAKS, ch


def test_void_thresholds_need_both_a_gap_and_an_empty_box():
    """The template's own slides leave the same band under the title and look
    right, because their copy fills the boxes - so the gap alone cannot be the
    rule. Measured: reference decks 0.44-0.58 fill, model-built ones 0.14-0.31."""
    assert A.VOID_GAP <= 0.12
    assert 0.18 < A.VOID_FILL < 0.44


class _W:
    """A pdftotext word box: only the vertical extent matters here."""

    def __init__(self, top, bottom):
        self.top, self.bottom = top, bottom


def _blank_slide():
    prs = Presentation()
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_rendered_void_flags_copy_that_stops_halfway_down():
    prs, slide = _blank_slide()
    h = prs.slide_height
    words = [_W(int(0.05 * h), int(0.30 * h))] * A.RENDERED_VOID_MIN_WORDS
    band = A.rendered_void(slide, words, prs.slide_width, h)
    assert band is not None and band > 0.65


def test_rendered_void_ignores_copy_that_runs_to_the_bottom():
    prs, slide = _blank_slide()
    h = prs.slide_height
    words = [_W(int(0.05 * h), int(0.20 * h)), _W(int(0.25 * h), int(0.95 * h))]
    words *= A.RENDERED_VOID_MIN_WORDS
    assert A.rendered_void(slide, words, prs.slide_width, h) is None


def test_rendered_void_spares_a_cover_slide():
    """A cover is a headline in the middle of an empty slide. Below the word
    floor the band is the design, not a hole."""
    prs, slide = _blank_slide()
    h = prs.slide_height
    words = [_W(int(0.40 * h), int(0.55 * h))] * 3
    assert A.rendered_void(slide, words, prs.slide_width, h) is None


def test_rendered_void_band_leaves_room_above_the_reference_decks():
    """Both AFTER_ decks and the two Dust templates top out at 0.46; the four
    Luna runs that shipped a hole measured 0.56-0.66."""
    assert 0.46 < A.RENDERED_VOID_BAND < 0.56


if __name__ == "__main__":
    tests = [v for k, v in sorted(globals().items())
             if k.startswith("test_") and callable(v)]
    for fn in tests:
        fn()
        print(f"ok   {fn.__name__}")
    print(f"\n{len(tests)} audit tests passed")
