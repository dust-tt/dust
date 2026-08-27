"""Tier-2 tests for pptx_slides.op_delete: a deleted slide leaves no part
behind, custom shows included. Builds decks with python-pptx, no soffice. Run
directly (`python test_slides.py`) or under pytest.

Lives in soffice/tests/, a subdir getLocalDirContent skips: it copies only the
regular files directly in soffice/ (never recursing), so tests never ship in the image. It adds soffice/ to sys.path to import the module.
"""
import sys
import tempfile
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from lxml import etree  # noqa: E402
from pptx import Presentation  # noqa: E402
from PIL import Image  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

import pptx_slides as S  # noqa: E402

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _deck(path: Path, slides: int = 3, custom_show_on: int = 0) -> None:
    """A deck; custom_show_on=N adds a custom show referencing slide N."""
    prs = Presentation()
    for _ in range(slides):
        prs.slides.add_slide(prs.slide_layouts[6])
    if custom_show_on:
        rid = list(prs.slides._sldIdLst)[custom_show_on - 1].get(qn("r:id"))
        cust = etree.SubElement(prs.part._element, f"{{{P_NS}}}custShowLst")
        show = etree.SubElement(cust, f"{{{P_NS}}}custShow")
        show.set("name", "Short version")
        show.set("id", "0")
        ref = etree.SubElement(
            etree.SubElement(show, f"{{{P_NS}}}sldLst"), f"{{{P_NS}}}sld"
        )
        ref.set(f"{{{R_NS}}}id", rid)
    prs.save(path)


def _slide_parts(path: Path) -> int:
    with zipfile.ZipFile(path) as zf:
        return len(
            [n for n in zf.namelist() if n.startswith("ppt/slides/slide")
             and n.endswith(".xml")]
        )


def _delete(path: Path, slide_nos):
    prs = Presentation(path)
    msg = S.op_delete(prs, slide_nos)
    prs.save(path)
    return msg


def test_delete_drops_the_slide_part():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck)
        _delete(deck, [2])
        assert len(Presentation(deck).slides) == 2
        assert _slide_parts(deck) == 2


def test_delete_drops_the_part_when_a_custom_show_also_references_it():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck, custom_show_on=2)
        _delete(deck, [2])
        assert len(Presentation(deck).slides) == 2
        # Without the purge, drop_rel no-ops and slide2.xml survives.
        assert _slide_parts(deck) == 2


def test_delete_leaves_the_custom_shows_other_slides_alone():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck, custom_show_on=2)
        _delete(deck, [3])
        assert len(Presentation(deck).slides) == 2
        assert _slide_parts(deck) == 2
        with zipfile.ZipFile(deck) as zf:
            pres = zf.read("ppt/presentation.xml").decode("utf-8")
        assert "custShow" in pres


def test_delete_out_of_range_raises():
    with tempfile.TemporaryDirectory() as d:
        deck = Path(d) / "deck.pptx"
        _deck(deck)
        for bad in ([0], [4], [1, 9]):
            try:
                _delete(deck, bad)
            except ValueError:
                continue
            raise AssertionError(f"expected ValueError for {bad!r}")

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def _deck_with_picture_background(path: Path, bg_image: Path, other: Path) -> None:
    """One slide with a notes slide, a picture background, and a picture shape.

    The notes slide is the point: pptx_slides skips it when copying rels, so the
    copy's relationship ids shift down by one against the source's. That shift is
    what turns an unremapped background r:embed into a pointer at the WRONG media
    part instead of a harmlessly identical one.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "speaker notes"
    _, bg_rid = slide.part.get_or_add_image_part(str(bg_image))
    slide.shapes.add_picture(str(other), 0, 0, 914400, 914400)

    bg = etree.Element(f"{{{P_NS}}}bg")
    bg_pr = etree.SubElement(bg, f"{{{P_NS}}}bgPr")
    blip_fill = etree.SubElement(bg_pr, f"{{{A_NS}}}blipFill")
    blip = etree.SubElement(blip_fill, f"{{{A_NS}}}blip")
    blip.set(f"{{{R_NS}}}embed", bg_rid)
    etree.SubElement(bg_pr, f"{{{A_NS}}}effectLst")
    slide._element.cSld.insert(0, bg)
    prs.save(path)


def _background_target(path: Path, slide_no: int):
    prs = Presentation(path)
    slide = list(prs.slides)[slide_no - 1]
    bg = slide._element.cSld.find(qn("p:bg"))
    if bg is None:
        return None
    blips = bg.findall(".//" + qn("a:blip"))
    if not blips:
        return None
    rid = blips[0].get(qn("r:embed"))
    try:
        return str(slide.part.rels[rid].target_partname)
    except KeyError:
        return "UNRESOLVED"


def test_duplicate_keeps_the_picture_background_pointing_at_its_own_image():
    """A duplicated slide's <p:bg> lives outside <p:spTree>, so rewriting only
    the shape tree left its r:embed on the SOURCE slide's relationship id - and
    the copy silently rendered a different image behind it."""
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        # Two distinct images, so a background pointing at the wrong one is
        # visible as a different media part rather than the same bytes twice.
        bg_image, other = tmp / "bg.png", tmp / "other.png"
        Image.new("RGB", (4, 4), (10, 20, 30)).save(bg_image)
        Image.new("RGB", (8, 8), (200, 100, 50)).save(other)
        deck = tmp / "d.pptx"
        _deck_with_picture_background(deck, bg_image, other)
        original = _background_target(deck, 1)
        assert original and original != "UNRESOLVED", original

        prs = Presentation(deck)
        S.op_duplicate(prs, [1], 1, 1)
        prs.save(deck)

        assert _background_target(deck, 2) == original, (
            f"copy's background resolves to {_background_target(deck, 2)}, "
            f"expected {original}"
        )


if __name__ == "__main__":
    tests = [v for k, v in sorted(globals().items())
             if k.startswith("test_") and callable(v)]
    for fn in tests:
        fn()
        print(f"ok   {fn.__name__}")
    print(f"\n{len(tests)} slides tests passed")
